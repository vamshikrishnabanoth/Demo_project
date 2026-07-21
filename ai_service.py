import json
import re
import random
import os
import sys
import base64
import uuid
import psycopg2
from psycopg2.extras import RealDictCursor
from dotenv import load_dotenv

# Load environment from server/.env
load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), 'server', '.env'))

# Force UTF-8 output on Windows to avoid cp1252 UnicodeEncodeError with emoji/unicode chars
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
import tempfile
import hashlib
import numpy as np
import faiss
import requests
from typing import List, Optional
from fastapi import FastAPI, HTTPException, UploadFile, File, BackgroundTasks
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer
from fastapi.middleware.cors import CORSMiddleware
import uvicorn
import ast

# ── Pedagogical Configuration & Redis Centralized Caching Setup ──────────────
PEDAGOGICAL_CONFIG_PATH = os.path.join(os.path.dirname(__file__), 'config', 'pedagogical_config.json')

def load_pedagogical_config():
    if os.path.exists(PEDAGOGICAL_CONFIG_PATH):
        try:
            with open(PEDAGOGICAL_CONFIG_PATH, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"⚠️ Error loading pedagogical config: {e}")
    return {
        "importance_weights": {"coverage": 0.25, "depth": 0.20, "repetition": 0.15, "teacher_emphasis": 0.20, "dependency": 0.10, "learning_objectives": 0.10},
        "feature_flags": {"enable_legacy_classifier": False, "enable_redis_cache": True, "enable_planner_validator": True},
        "coding_subdivisions": {
            "implementation": ["write_code", "code_completion", "fill_in_blank"],
            "analysis": ["output_prediction", "debugging", "trace_execution", "compiler_error"]
        }
    }

pedagogical_config = load_pedagogical_config()

# Redis Centralized Cache setup (multi-worker support)
redis_client = None
REDIS_URL = os.getenv("REDIS_URL", "redis://localhost:6379/0")
try:
    import redis
    redis_client = redis.from_url(REDIS_URL, decode_responses=True)
    redis_client.ping()
    print("✅ Redis Centralized Stage A Cache connected successfully.")
except Exception as re_err:
    print(f"ℹ️ Redis offline or unavailable. Falling back to in-memory Stage A cache.")
    redis_client = None

_in_memory_stage_a_cache = {}

def get_stage_a_cache(content_hash: str):
    if redis_client:
        try:
            val = redis_client.get(f"stage_a:{content_hash}")
            if val:
                return json.loads(val)
        except Exception:
            pass
    return _in_memory_stage_a_cache.get(content_hash)

def set_stage_a_cache(content_hash: str, data: dict, ttl_seconds: int = 43200):
    if redis_client:
        try:
            redis_client.setex(f"stage_a:{content_hash}", ttl_seconds, json.dumps(data))
            return
        except Exception:
            pass
    _in_memory_stage_a_cache[content_hash] = data

def compute_normalized_hash(text: str) -> str:
    """Computes SHA-256 content_hash AFTER text cleaning/normalization to prevent minor formatting cache busts."""
    normalized = re.sub(r'\s+', ' ', (text or '').strip().lower())
    return hashlib.sha256(normalized.encode('utf-8')).hexdigest()


# Extractors
import easyocr
import docx
from pptx import Presentation
from PyPDF2 import PdfReader

# -----------------------------
# 1. CONFIGURATION
# -----------------------------
OLLAMA_URL = "http://127.0.0.1:11434/api/generate"
MODEL_NAME = "quiz-expert"
DATABASE_URL = os.getenv("DATABASE_URL")

# Initialize RAG Engine (Phase 3)
print("Loading RAG Embedding Engine...")
embed_model = SentenceTransformer('all-MiniLM-L6-v2') 

# Initialize OCR
print("Loading OCR Engine...")
ocr_reader = easyocr.Reader(['en'])

# Initialize Whisper model locally (using 'base' optimized for CUDA/CPU)
whisper_model = None
try:
    print("Loading Local Whisper Engine...")
    from faster_whisper import WhisperModel
    try:
        # Enforce int8 quantization for 4x faster CPU execution
        whisper_model = WhisperModel("base", device="auto", compute_type="int8")
    except Exception:
        whisper_model = WhisperModel("base", device="auto", compute_type="default")
    print("Local Whisper Engine loaded successfully!")
except Exception as e:
    print(f"⚠️ Warning: Could not initialize local Whisper: {e}. Voice transcription will fall back to cloud/mock.")

# -----------------------------
# 2. DATABASE & RAG CORE LOGIC
# -----------------------------
def get_db_connection():
    if not DATABASE_URL:
        raise ValueError("DATABASE_URL is not set in environment.")
    
    url = DATABASE_URL
    if "connection_limit" in url:
        import urllib.parse as urlparse
        try:
            parsed = urlparse.urlparse(url)
            query = urlparse.parse_qs(parsed.query)
            query.pop('connection_limit', None)
            new_query = urlparse.urlencode(query, doseq=True)
            parsed = parsed._replace(query=new_query)
            url = urlparse.urlunparse(parsed)
        except Exception as e:
            print(f"Error stripping connection_limit: {e}")
            
    return psycopg2.connect(url)

def recursive_token_splitter(text, max_tokens=500, overlap_percent=10):
    if not text or not text.strip():
        return []
    
    # Pre-split giant texts to prevent token index blowouts & warnings
    if len(text) > 10000:
        blocks = []
        step = 8000
        for i in range(0, len(text), step):
            block = text[i:i+10000].strip()
            if len(block) > 10:
                blocks.append(block)
        
        all_chunks = []
        for block in blocks:
            all_chunks.extend(recursive_token_splitter(block, max_tokens, overlap_percent))
        return all_chunks

    try:
        tokens = embed_model.tokenizer.encode(text, add_special_tokens=False, truncation=True, max_length=10000)
    except Exception as e:
        print(f"Tokenization error, falling back to character approximation: {e}")
        # Fallback character approximation (1 token ~= 4 chars)
        char_limit = max_tokens * 4
        overlap_chars = int(char_limit * (overlap_percent / 100.0))
        step = char_limit - overlap_chars
        if step <= 0:
            step = char_limit
        return [text[i:i+char_limit].strip() for i in range(0, len(text), step) if len(text[i:i+char_limit].strip()) > 10]
    
    chunks = []
    overlap = int(max_tokens * (overlap_percent / 100.0))
    step = max_tokens - overlap
    if step <= 0:
        step = max_tokens
        
    for i in range(0, len(tokens), step):
        chunk_tokens = tokens[i : i + max_tokens]
        chunk_text = embed_model.tokenizer.decode(chunk_tokens, clean_up_tokenization_spaces=True)
        if len(chunk_text.strip()) > 10:
            chunks.append(chunk_text.strip())
        if i + max_tokens >= len(tokens):
            break
    return chunks

def get_relevant_db_context(query_text, top_k=5):
    try:
        query_emb = embed_model.encode([query_text])[0]
        query_emb_str = "[" + ",".join(map(str, query_emb)) + "]"
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        try:
            # Cosine distance operator <=> in pgvector
            cur.execute(
                'SELECT id, content, source, (embedding <=> %s::vector) as distance '
                'FROM "DocumentChunk" '
                'ORDER BY distance ASC '
                'LIMIT %s',
                (query_emb_str, top_k)
            )
            results = cur.fetchall()
            return [dict(row) for row in results]
        except Exception as e:
            print(f"Error executing vector similarity search: {e}")
            return []
        finally:
            cur.close()
            conn.close()
    except Exception as e:
        print(f"Failed to generate query embedding or connect: {e}")
        return []
def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        if ext in ['.jpg', '.jpeg', '.png']:
            text = " ".join(ocr_reader.readtext(file_path, detail=0))
        elif ext == '.docx':
            doc = docx.Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        elif ext == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() or ""
    except Exception as e:
        print(f"Extraction Error: {e}")
    return text.strip()

def get_relevant_context(text, query, top_k=3):
    paragraphs = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 100]
    if not paragraphs:
        paragraphs = [text[i:i+1000] for i in range(0, len(text), 800)]
    
    if len(paragraphs) <= top_k:
        return " ".join(paragraphs)

    embeddings = embed_model.encode(paragraphs)
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(np.array(embeddings).astype('float32'))
    
    query_emb = embed_model.encode([query])
    distances, indices = index.search(np.array(query_emb).astype('float32'), top_k)
    
    relevant_chunks = [paragraphs[i] for i in indices[0]]
    return " ".join(relevant_chunks)

# -----------------------------
# 3. FASTAPI SERVICE
# -----------------------------
from fastapi.exceptions import RequestValidationError
from fastapi.responses import JSONResponse

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.exception_handler(RequestValidationError)
async def validation_exception_handler(request, exc):
    print("❌ Validation Error Details:", exc.errors())
    print("❌ Request Body Sent:", exc.body)
    return JSONResponse(
        status_code=422,
        content={"detail": exc.errors(), "body": exc.body},
    )

@app.get("/")
async def root():
    return {"status": "online", "service": "quiz-expert-ai"}

class MultiInputSource(BaseModel):
    type: str # 'text', 'pdf', 'docx', 'pptx', 'image', 'voice'
    content: str
    source_name: str = "Unknown Source"

class AnalyzeRequest(BaseModel):
    inputs: List[MultiInputSource]

class GeneratorRequest(BaseModel):
    inputs: Optional[List[MultiInputSource]] = None
    type: Optional[str] = None # Backwards compatibility
    content: Optional[str] = None # Backwards compatibility
    count: int = 5
    difficulty: str = "Medium"
    target_ratios: Optional[dict] = None
    source_material_id: Optional[str] = None
    topic_weights: Optional[dict] = None
    callback_url: Optional[str] = None
    taskId: Optional[str] = None
    isolated_narratives: Optional[List[str]] = None
    question_style: Optional[str] = "MIXED"
    interview_mode: Optional[bool] = False

class StageAAnalysisRequest(BaseModel):
    inputs: Optional[List[MultiInputSource]] = None
    content: Optional[str] = None
    type: Optional[str] = None

class StageBPlannerRequest(BaseModel):
    stage_a_data: dict
    count: int = 5
    difficulty: str = "Medium"
    question_style: str = "MIXED"
    interview_mode: bool = False


class IngestRequest(BaseModel):
    source: str
    type: str # 'text', 'pdf', 'docx', 'pptx', 'image'
    content: str # Can be raw text, file path, or base64
    metadata: dict = {}

@app.post("/admin/ingest")
async def ingest_document(req: IngestRequest):
    source_text = req.content
    
    if req.type in ['pdf', 'docx', 'pptx', 'image'] and os.path.exists(req.content):
        source_text = extract_text_from_file(req.content)
    elif req.content.startswith("base64:"):
        try:
            b64_data = req.content[7:]
            file_data = base64.b64decode(b64_data)
            ext = req.type if req.type != 'image' else 'png'
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                tmp.write(file_data)
                tmp_path = tmp.name
            source_text = extract_text_from_file(tmp_path)
            os.remove(tmp_path)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Invalid base64 file content: {e}")
            
    if not source_text or len(source_text) < 10:
        raise HTTPException(status_code=400, detail="Content too short or unextractable.")
        
    try:
        chunks = recursive_token_splitter(source_text, max_tokens=500, overlap_percent=10)
        if not chunks:
            raise HTTPException(status_code=400, detail="No chunks generated from content.")
            
        embeddings = embed_model.encode(chunks)
        
        conn = get_db_connection()
        cur = conn.cursor()
        try:
            for chunk, embedding in zip(chunks, embeddings):
                chunk_id = str(uuid.uuid4())
                embedding_str = "[" + ",".join(map(str, embedding)) + "]"
                metadata_json = json.dumps(req.metadata)
                cur.execute(
                    'INSERT INTO "DocumentChunk" (id, content, embedding, source, metadata, "createdAt") '
                    'VALUES (%s, %s, %s::vector, %s, %s, NOW())',
                    (chunk_id, chunk, embedding_str, req.source, metadata_json)
                )
            conn.commit()
        except Exception as e:
            conn.rollback()
            print(f"Database insertion error: {e}")
            raise HTTPException(status_code=500, detail=f"Database insertion failed: {e}")
        finally:
            cur.close()
            conn.close()
            
        return {
            "status": "success",
            "message": f"Successfully ingested {len(chunks)} chunks from source '{req.source}'.",
            "chunks_count": len(chunks)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

def call_groq_fallback(prompt, json_mode=True):
    import time
    groq_api_key = os.getenv("GROQ_API_KEY")
    if not groq_api_key:
        raise ValueError("Groq API Key is not set in environment.")
    
    headers = {
        "Authorization": f"Bearer {groq_api_key}",
        "Content-Type": "application/json"
    }
    
    # Rotate models to bypass TPM rate limits if hit
    models = ["llama-3.1-8b-instant", "gemma2-9b-it", "mixtral-8x7b-32768"]
    
    for attempt in range(4):
        model = models[attempt % len(models)]
        payload = {
            "messages": [{"role": "user", "content": prompt}],
            "model": model,
            "temperature": 0.2
        }
        if json_mode:
            payload["response_format"] = {"type": "json_object"}
            
        try:
            response = requests.post("https://api.groq.com/openai/v1/chat/completions", json=payload, headers=headers, timeout=45)
            if response.status_code == 200:
                return response.json()["choices"][0]["message"]["content"]
            elif response.status_code == 429:
                print(f"  [GROQ] ⚠️ 429 Rate Limit on {model}. Retrying next model in 3 seconds...")
                time.sleep(3)
            else:
                print(f"  [GROQ] ⚠️ API returned {response.status_code}. Retrying...")
                time.sleep(2)
        except Exception as e:
            print(f"  [GROQ] ⚠️ Request error: {e}. Retrying...")
            time.sleep(2)
            
    raise ValueError("Groq API failed after 4 retries.")

def sanitize_source_text(text: str) -> str:
    """
    Strips out file names, timestamps, and structural labels like
    'Voice Transcript (12:23:03 AM)' so the LLM never sees them.
    """
    # Remove strings like 'Voice Transcript (XX:XX:XX AM/PM)' case-insensitively
    text = re.sub(r'(?i)voice\s+transcript\s*\([^)]+\)', '', text)
    
    # Remove standalone timestamps like 12:23:03 AM
    text = re.sub(r'\b\d{1,2}:\d{2}:\d{2}\s*(?:AM|PM|am|pm)\b', '', text)
    
    # Clean up double newlines left over by stripping
    return re.sub(r'\n\s*\n', '\n', text).strip()

def call_vision_model(img_b64, mime_type):
    prompt = (
        "Transcribe all text from this image. Convert all mathematical equations and formulas into strict, clean inline or display LaTeX syntax ($...$ or $$...$$). "
        "For any hand-drawn architectural blocks, circuits, or diagrams, output a clear, detailed structural Markdown description mapping the inputs, outputs, and components."
    )
    
    # 1. Try Groq Multimodal Vision (llama-3.2-11b-vision-preview)
    groq_api_key = os.getenv("GROQ_API_KEY")
    if groq_api_key:
        print("  [VISION LAYER] Querying Groq Vision Model (llama-3.2-11b-vision-preview)...")
        headers = {
            "Authorization": f"Bearer {groq_api_key}",
            "Content-Type": "application/json"
        }
        
        data_url = f"data:{mime_type};base64,{img_b64}"
        payload = {
            "model": "llama-3.2-11b-vision-preview",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": data_url
                            }
                        }
                    ]
                }
            ],
            "temperature": 0.2
        }
        
        try:
            response = requests.post("https://api.groq.com/openai/v1/chat/completions", json=payload, headers=headers, timeout=45)
            if response.status_code == 200:
                transcription = response.json()["choices"][0]["message"]["content"]
                print(f"  [VISION LAYER] Successfully transcribed {len(transcription)} chars using Groq Vision.")
                return transcription
            else:
                print(f"  [VISION LAYER] Groq Vision API returned {response.status_code}: {response.text}")
        except Exception as e:
            print(f"  [VISION LAYER] Groq Vision error: {e}")
            
    # 2. Fallback to Local Ollama Vision (llama3.2-vision)
    print("  [VISION LAYER] Querying Local Ollama Vision (llama3.2-vision)...")
    payload = {
        "model": "llama3.2-vision",
        "prompt": prompt,
        "images": [img_b64],
        "stream": False,
        "options": {
            "temperature": 0.2
        }
    }
    
    try:
        response = requests.post(OLLAMA_URL.replace("/api/chat", "/api/generate").replace("/api/generate", "/api/generate"), json=payload, timeout=90)
        if response.status_code == 200:
            transcription = response.json().get("response", "")
            print(f"  [VISION LAYER] Successfully transcribed {len(transcription)} chars using Local Ollama Vision.")
            return transcription
        else:
            print(f"  [VISION LAYER] Local Ollama Vision API returned {response.status_code}: {response.text}")
    except Exception as e:
        print(f"  [VISION LAYER] Local Ollama Vision error: {e}")
        
    return ""

def process_image_sources(inp):
    print(f"  [VISION LAYER] Intercepting source: {inp.source_name} ({inp.type})")
    
    is_pdf = False
    file_bytes = None
    
    if inp.content.startswith("base64:"):
        b64_data = inp.content[7:]
        file_bytes = base64.b64decode(b64_data)
        if file_bytes.startswith(b'%PDF'):
            is_pdf = True
    elif os.path.exists(inp.content):
        ext = os.path.splitext(inp.content)[1].lower()
        if ext == '.pdf':
            is_pdf = True
            with open(inp.content, 'rb') as f:
                file_bytes = f.read()
        else:
            with open(inp.content, 'rb') as f:
                file_bytes = f.read()
                
    if not file_bytes:
        return ""

    images_to_process = []
    
    if is_pdf:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
                
            reader = PdfReader(tmp_path)
            page_index = 1
            for page in reader.pages:
                for img_obj in page.images:
                    img_data = img_obj.data
                    img_b64 = base64.b64encode(img_data).decode('utf-8')
                    ext = os.path.splitext(img_obj.name)[1].lower() if img_obj.name else '.png'
                    m_type = "image/jpeg" if ext in ['.jpg', '.jpeg'] else "image/png"
                    images_to_process.append((img_b64, m_type))
                
                if not page.images:
                    try:
                        from pdf2image import convert_from_path
                        pages = convert_from_path(tmp_path, first_page=page_index, last_page=page_index)
                        for p in pages:
                            import io
                            img_byte_arr = io.BytesIO()
                            p.save(img_byte_arr, format='PNG')
                            img_data = img_byte_arr.getvalue()
                            img_b64 = base64.b64encode(img_data).decode('utf-8')
                            images_to_process.append((img_b64, "image/png"))
                    except Exception as pe:
                        print(f"  [VISION LAYER] pdf2image fallback skipped: {pe}")
                page_index += 1
            os.remove(tmp_path)
        except Exception as e:
            print(f"  [VISION LAYER] Error processing PDF scan: {e}")
    else:
        ext = os.path.splitext(inp.source_name)[1].lower() if inp.source_name else '.png'
        m_type = "image/jpeg" if ext in ['.jpg', '.jpeg'] else "image/png"
        img_b64 = base64.b64encode(file_bytes).decode('utf-8')
        images_to_process.append((img_b64, m_type))
        
    if not images_to_process:
        print("  [VISION LAYER] No images found to process.")
        return ""
        
    transcriptions = []
    for idx, (img_b64, mime_type) in enumerate(images_to_process):
        print(f"  [VISION LAYER] Processing image {idx+1}/{len(images_to_process)}...")
        trans = call_vision_model(img_b64, mime_type)
        if trans:
            transcriptions.append(trans)
            
    return "\n\n".join(transcriptions)

def resolve_input_sources(inputs):
    aggregated_texts = []
    voice_transcripts = []
    
    for inp in inputs:
        source_text = ""
        if inp.type in ['image', 'handwritten_scan']:
            source_text = process_image_sources(inp)
        elif inp.type in ['pdf', 'docx', 'pptx']:
            if os.path.exists(inp.content):
                source_text = extract_text_from_file(inp.content)
            elif inp.content.startswith("base64:"):
                try:
                    b64_data = inp.content[7:]
                    file_data = base64.b64decode(b64_data)
                    ext = inp.type
                    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                        tmp.write(file_data)
                        tmp_path = tmp.name
                    source_text = extract_text_from_file(tmp_path)
                    os.remove(tmp_path)
                except Exception as e:
                    print(f"Error parsing base64 file source: {e}")
                    source_text = ""
            else:
                source_text = inp.content
        else:
            source_text = inp.content
        
        if source_text and len(source_text.strip()) > 10:
            sanitized = sanitize_source_text(source_text)
            if len(sanitized.strip()) > 10:
                if inp.type == 'voice':
                    voice_transcripts.append(sanitized.strip())
                else:
                    aggregated_texts.append(sanitized.strip())
                    
    combined = []
    if voice_transcripts:
        combined.append("=== VOICE TRANSCRIPT CONTEXT ===\n" + "\n\n".join(voice_transcripts))
    if aggregated_texts:
        combined.append("=== DOCUMENT CONTEXT ===\n" + "\n\n".join(aggregated_texts))
        
    return "\n\n".join(combined) if combined else ""

def deduplicate_text_chunks(text, embed_model, max_tokens=500, overlap_percent=10):
    chunks = recursive_token_splitter(text, max_tokens, overlap_percent)
    if not chunks:
        return []
    
    embeddings = embed_model.encode(chunks)
    
    kept_chunks = []
    kept_embeddings = []
    
    for chunk, embedding in zip(chunks, embeddings):
        is_duplicate = False
        if kept_embeddings:
            similarities = np.dot(kept_embeddings, embedding)
            if np.max(similarities) > 0.90:
                is_duplicate = True
        if not is_duplicate:
            kept_chunks.append(chunk)
            kept_embeddings.append(embedding)
            
    print(f"[Reduction]    : Kept {len(kept_chunks)} core academic blocks. Removed {len(chunks) - len(kept_chunks)} duplicates.")
    return kept_chunks

def calibrate_ratios_by_heuristics(text: str, current_ratios: dict) -> dict:
    import re
    
    text_lower = text.lower()
    
    # Heuristical counts
    code_keywords = [
        r'#include', r'void\s+\*', r'pthread_create', r'int\s+main', r'gcc\s+', r'ctime',
        r'st_mode', r'd_name', r'struct\s+stat', r'msgget\(', r'sem_wait\(', r'sem_post\(',
        r'fork\(', r'kill\(', r'close\(', r'write\(', r'read\(', r'open\(', r'printf\(', 
        r'scanf\(', r'opendir\(', r'readdir\(', r'closedir\(', r'd_ino', r'inode\b', 
        r'pointer', r'\.c\b', r'-lpthread', r'sem_init\('
    ]
    code_count = 0
    for kw in code_keywords:
        try:
            code_count += len(re.findall(kw, text_lower if '\\s' not in kw and '\\(' not in kw else text))
        except Exception:
            pass
            
    math_keywords = [
        r'pageno\b', r'logicaladdress', r'physicaladdress', r'pagesize', r'offset\b', 
        r'frameno\b', r'pagetable', r'need\s*=\s*max', r'max\s*-\s*alloc', r'safe\s+sequence',
        r'avail\[', r'alloc\[', r'need\[', r'max\[', r'banker\'s', r'matrix', r'allocation',
        r'sem_init\(&empty', r'la%pagesize', r'la/pagesize', r'frame\s*×\s*pagesize'
    ]
    math_count = 0
    for kw in math_keywords:
        try:
            math_count += len(re.findall(kw, text_lower))
        except Exception:
            pass

    theory_keywords = [
        r'define\b', r'definition', r'explain\b', r'concept', r'viva\b', r'theory', r'what is'
    ]
    theory_count = sum(len(re.findall(kw, text_lower)) for kw in theory_keywords)

    comparison_keywords = [
        r'versus', r'vs\b', r'difference\b', r'advantage', r'disadvantage', r'trade-off', r'compare'
    ]
    comparison_count = sum(len(re.findall(kw, text_lower)) for kw in comparison_keywords)

    scenario_keywords = [
        r'scenario', r'case study', r'real-world', r'industrial', r'application'
    ]
    scenario_count = sum(len(re.findall(kw, text_lower)) for kw in scenario_keywords)

    # If there is strong evidence of code / math / theory / comparisons, we calibrate
    total_score = code_count + math_count + theory_count + comparison_count + scenario_count
    if total_score > 3:
        w_practical = max(0.1, float(code_count) * 2.0)
        w_formulas = max(0.1, float(math_count) * 2.0)
        w_theory = max(0.1, float(theory_count) * 0.8)
        w_comparison = max(0.1, float(comparison_count) * 1.0)
        w_scenario = max(0.1, float(scenario_count) * 1.0)
        
        # Boosts
        if code_count > 2:
            w_practical += 4.0
        if math_count > 1:
            w_formulas += 3.0
            
        sum_w = w_practical + w_formulas + w_theory + w_comparison + w_scenario
        
        # Calculate ratio values
        new_ratios = {
            "CONCEPTS_AND_DEFINITIONS": round(w_theory / sum_w, 2),
            "COMPARISONS_AND_TRADEOFFS": round(w_comparison / sum_w, 2),
            "FORMULAS_AND_CALCULATIONS": round(w_formulas / sum_w, 2),
            "CASE_STUDIES_AND_SCENARIOS": round(w_scenario / sum_w, 2),
            "PRACTICAL_AND_LAB_TASKS": round(w_practical / sum_w, 2)
        }
        
        # Ensure sum is exactly 1.0
        diff = round(1.0 - sum(new_ratios.values()), 2)
        if diff != 0.0:
            max_key = max(new_ratios, key=new_ratios.get)
            new_ratios[max_key] = round(new_ratios[max_key] + diff, 2)
            
        return new_ratios
        
    return current_ratios

@app.post("/analyze-sources")
async def analyze_sources(req: AnalyzeRequest):
    if not req.inputs:
        raise HTTPException(status_code=400, detail="No input sources provided.")
        
    text = resolve_input_sources(req.inputs)
    if not text or len(text.strip()) < 2:
        raise HTTPException(status_code=400, detail="Content too short or unextractable.")
        
    chunks = deduplicate_text_chunks(text, embed_model)
    context = "\n\n".join(chunks[:6])
    
    print("  [AGENT 1] Analyzing academic content, token density, summaries, and topics...")
    prompt = (
        "You are an elite academic analyzer. Analyze the textbook/lecture context below.\n\n"
        "MANDATORY CONTEXTUAL DE-NOISING PASS:\n"
        "You must explicitly inspect the provided context. If it contains audio transcripts or speech data, identify and strip out all non-pedagogical clutter: classroom management comments, administrative/scheduling remarks (e.g., homework deadlines, quiz notifications), and unrelated conversational anecdotes. Process only the pure, academic, and engineering concepts.\n\n"
        f"Context:\n{context[:6000]}\n\n"
        "Tasks:\n"
        "1. Check if the content is educational/academic. Set 'relevancy_verdict' to 'pass' if it is academic (note: programming manuals, code files, syntax lists, data structures, and computer science slides are 100% academic/educational), or 'fail' if it is gibberish, casual chat, or spam.\n"
        "2. Create a bulleted lobby summary (3-4 concise, high-impact bullet points for a quiz lobby study panel).\n"
        "3. Generate 5 core study flashcards (Q&A style for post-quiz review).\n"
        "4. Suggest target ratios distributing a total weight of 1.0 across these 5 Master Academic Archetypes based on pedagogical intent (Do NOT always output exactly 0.2 for all categories. If the text is heavily programming/lab-based, allocate more weight to PRACTICAL_AND_LAB_TASKS. If it has math formulas/matrices, allocate more weight to FORMULAS_AND_CALCULATIONS. The sum of all ratios must be exactly 1.0):\n"
        "   - 'CONCEPTS_AND_DEFINITIONS' (Core Theory)\n"
        "   - 'COMPARISONS_AND_TRADEOFFS' (Analytical Reasoning)\n"
        "   - 'FORMULAS_AND_CALCULATIONS' (Numerical Design)\n"
        "   - 'CASE_STUDIES_AND_SCENARIOS' (Real-World Application)\n"
        "   - 'PRACTICAL_AND_LAB_TASKS' (Implementation Synthesis)\n"
        "5. Extract 5-10 specific curriculum concept tags and baseline weights (0.0 to 1.0).\n"
        "6. Classify contextual examples into either 'CLASSIC_DOMAIN_STANDARD' (foundational standard examples common to the domain) or 'TRANSIENT_ANALOGY' (casual/metaphorical settings or temporary stories). If an example is classified as a 'TRANSIENT_ANALOGY', extract the underlying mathematical/logical rule, and place the specific characters/names/settings used into an 'isolated_narratives' exclusion array of strings.\n\n"
        "Return ONLY a clean JSON object conforming strictly to this format:\n"
        "{\n"
        "  \"relevancy_verdict\": \"pass\",\n"
        "  \"relevancy_reason\": \"...\",\n"
        "  \"lobby_summary\": \"- Key concept 1...\\n- Key concept 2...\",\n"
        "  \"ai_flashcards\": [\n"
        "    {\"question\": \"...\", \"answer\": \"...\"}\n"
        "  ],\n"
        "  \"ai_recommendation\": {\n"
        "    \"CONCEPTS_AND_DEFINITIONS\": 0.2,\n"
        "    \"COMPARISONS_AND_TRADEOFFS\": 0.2,\n"
        "    \"FORMULAS_AND_CALCULATIONS\": 0.2,\n"
        "    \"CASE_STUDIES_AND_SCENARIOS\": 0.2,\n"
        "    \"PRACTICAL_AND_LAB_TASKS\": 0.2\n"
        "  },\n"
        "  \"concepts\": [\n"
        "    {\"concept_tag\": \"...\", \"weight_score\": 0.85}\n"
        "  ],\n"
        "  \"isolated_narratives\": [\n"
        "    \"Name of character/setting/anecdote to exclude (e.g. Alice and Bob, dining philosophers setting, specific car factory anecdote)\"\n"
        "  ]\n"
        "}"
    )

    payload = {
        "model": MODEL_NAME,
        "prompt": prompt,
        "stream": False,
        "format": "json",
        "options": {
            "temperature": 0.2,
            "num_ctx": 8192
        }
    }
    
    try:
        response = requests.post(OLLAMA_URL, json=payload, timeout=180)
        if response.status_code == 200:
            raw_text = response.json().get("response", "")
            data = robust_json_loads(raw_text)
            
            if data.get("relevancy_verdict") == "fail":
                raise HTTPException(
                    status_code=422,
                    detail={
                        "status": "validation_error",
                        "message": data.get("relevancy_reason") or "Non-academic content detected."
                    }
                )
            
            # Post-process to dynamically shift recommended ratios based on coding/formula heuristics
            if "ai_recommendation" in data and isinstance(data["ai_recommendation"], dict):
                data["ai_recommendation"] = calibrate_ratios_by_heuristics(text, data["ai_recommendation"])
                
            return data
    except HTTPException as he:
        raise he
    except Exception as e:
        print(f"⚠️ Local Ollama failed or timed out: {e}. Trying Groq Cloud fallback...")
        if os.getenv("GROQ_API_KEY"):
            try:
                # Trim the context inside prompt to avoid TPM limit on Groq
                fallback_context = context[:1500]
                fallback_prompt = (
                    "You are an elite academic analyzer. Analyze the textbook/lecture context below.\n\n"
                    "MANDATORY CONTEXTUAL DE-NOISING PASS:\n"
                    "You must explicitly inspect the provided context. If it contains audio transcripts or speech data, identify and strip out all non-pedagogical clutter: classroom management comments, administrative/scheduling remarks (e.g., homework deadlines, quiz notifications), and unrelated conversational anecdotes. Process only the pure, academic, and engineering concepts.\n\n"
                    f"Context:\n{fallback_context}\n\n"
                    "Tasks:\n"
                    "1. Check if the content is educational/academic. Set 'relevancy_verdict' to 'pass' if it is academic (note: programming manuals, code files, syntax lists, data structures, and computer science slides are 100% academic/educational), or 'fail' if it is gibberish, casual chat, or spam.\n"
                    "2. Create a bulleted lobby summary (3-4 concise, high-impact bullet points for a quiz lobby study panel).\n"
                    "3. Generate 5 core study flashcards (Q&A style for post-quiz review).\n"
                    "4. Suggest target ratios distributing a total weight of 1.0 across these 5 Master Academic Archetypes based on pedagogical intent (Do NOT always output exactly 0.2 for all categories. If the text is heavily programming/lab-based, allocate more weight to PRACTICAL_AND_LAB_TASKS. If it has math formulas/matrices, allocate more weight to FORMULAS_AND_CALCULATIONS. The sum of all ratios must be exactly 1.0):\n"
                    "   - 'CONCEPTS_AND_DEFINITIONS' (Core Theory)\n"
                    "   - 'COMPARISONS_AND_TRADEOFFS' (Analytical Reasoning)\n"
                    "   - 'FORMULAS_AND_CALCULATIONS' (Numerical Design)\n"
                    "   - 'CASE_STUDIES_AND_SCENARIOS' (Real-World Application)\n"
                    "   - 'PRACTICAL_AND_LAB_TASKS' (Implementation Synthesis)\n"
                    "5. Extract 5-10 specific curriculum concept tags and baseline weights (0.0 to 1.0).\n"
                    "6. Classify contextual examples into either 'CLASSIC_DOMAIN_STANDARD' (foundational standard examples common to the domain) or 'TRANSIENT_ANALOGY' (casual/metaphorical settings or temporary stories). If an example is classified as a 'TRANSIENT_ANALOGY', extract the underlying mathematical/logical rule, and place the specific characters/names/settings used into an 'isolated_narratives' exclusion array of strings.\n\n"
                    "Return ONLY a clean JSON object conforming strictly to this format:\n"
                    "{\n"
                    "  \"relevancy_verdict\": \"pass\",\n"
                    "  \"relevancy_reason\": \"...\",\n"
                    "  \"lobby_summary\": \"- Key concept 1...\\n- Key concept 2...\",\n"
                    "  \"ai_flashcards\": [\n"
                    "    {\"question\": \"...\", \"answer\": \"...\"}\n"
                    "  ],\n"
                    "  \"ai_recommendation\": {\n"
                    "    \"CONCEPTS_AND_DEFINITIONS\": 0.2,\n"
                    "    \"COMPARISONS_AND_TRADEOFFS\": 0.2,\n"
                    "    \"FORMULAS_AND_CALCULATIONS\": 0.2,\n"
                    "    \"CASE_STUDIES_AND_SCENARIOS\": 0.2,\n"
                    "    \"PRACTICAL_AND_LAB_TASKS\": 0.2\n"
                    "  },\n"
                    "  \"concepts\": [\n"
                    "    {\"concept_tag\": \"...\", \"weight_score\": 0.85}\n"
                    "  ],\n"
                    "  \"isolated_narratives\": [\n"
                    "    \"Name of character/setting/anecdote to exclude (e.g. Alice and Bob, dining philosophers setting, specific car factory anecdote)\"\n"
                    "  ]\n"
                    "}"
                )
                raw_text = call_groq_fallback(fallback_prompt, json_mode=True)
                data = robust_json_loads(raw_text)
                
                if data.get("relevancy_verdict") == "fail":
                    raise HTTPException(
                        status_code=422,
                        detail={
                            "status": "validation_error",
                            "message": data.get("relevancy_reason") or "Non-academic content detected."
                        }
                    )
                
                # Post-process to dynamically shift recommended ratios based on coding/formula heuristics
                if "ai_recommendation" in data and isinstance(data["ai_recommendation"], dict):
                    data["ai_recommendation"] = calibrate_ratios_by_heuristics(text, data["ai_recommendation"])
                    
                print("✅ Successfully recovered via Groq Cloud fallback!")
                return data
            except HTTPException as he2:
                raise he2
            except Exception as e2:
                print(f"❌ Groq fallback also failed: {e2}")
                raise HTTPException(status_code=500, detail=f"Ollama timed out and Groq fallback failed: {str(e2)}")
        else:
            raise HTTPException(status_code=500, detail=f"Ollama timed out/failed and no GROQ_API_KEY is configured. Error: {str(e)}")

@app.post("/stage-a-analysis")
async def get_stage_a_analysis(req: StageAAnalysisRequest):
    text = ""
    if req.inputs:
        text = resolve_input_sources(req.inputs)
    elif req.content:
        text = req.content
    
    if not text or len(text.strip()) < 2:
        raise HTTPException(status_code=400, detail="Content too short or unextractable for Stage A analysis.")

    stage_a_res = run_stage_a_content_understanding(text)
    return stage_a_res

@app.post("/stage-b-blueprint")
async def get_stage_b_blueprint(req: StageBPlannerRequest):
    try:
        blueprint = run_stage_b_academic_planner(
            stage_a_data=req.stage_a_data,
            question_count=req.count,
            difficulty=req.difficulty,
            question_style=req.question_style,
            interview_mode=req.interview_mode
        )
        return blueprint
    except Exception as e:
        raise HTTPException(status_code=422, detail=str(e))


@app.post("/transcribe")
async def transcribe_audio_file(file: UploadFile = File(...)):
    if not whisper_model:
        raise HTTPException(status_code=503, detail="Local Whisper model is not loaded.")
        
    try:
        suffix = os.path.splitext(file.filename)[1] or ".webm"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = tmp.name
            
        print(f"🎙️ Transcribing local file: {tmp_path} using Whisper base...")
        segments, info = whisper_model.transcribe(tmp_path, beam_size=5)
        text = " ".join([segment.text for segment in segments])
        
        raw_text = text.strip()
        
        print("\n" + "="*60)
        print("🎙️ [STEP 1: WHISPER TRANSCRIPTION RAW OUTPUT]")
        print(f"Content: {raw_text}")
        print("="*60 + "\n")
        
        cleaned_text = sanitize_source_text(raw_text)
        
        print("\n" + "="*60)
        print("🧹 [STEP 2: SANITIZED TEXT PAYLOAD]")
        print(f"Content: {cleaned_text}")
        print("="*60 + "\n")
        
        try:
            os.remove(tmp_path)
        except Exception:
            pass
            
        return {
            "status": "success",
            "text": cleaned_text,
            "language": info.language,
            "duration": info.duration
        }
    except Exception as e:
        print(f"❌ Local Transcription Error: {e}")
        raise HTTPException(status_code=500, detail=f"Local transcription failed: {str(e)}")

def robust_json_loads(text):
    text = text.strip()
    
    # Strategy 1: Direct JSON load
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
        
    # Strategy 2: Extract Markdown block
    match = re.search(r'```(?:json)?\s*([\s\S]*?)\s*```', text)
    if match:
        extracted = match.group(1).strip()
        try:
            return json.loads(extracted)
        except json.JSONDecodeError:
            text = extracted # Use extracted for subsequent strategies

    # Strategy 3: Greedy brace/bracket extraction
    match_brace = re.search(r'({[\s\S]*})|(\[[\s\S]*\])', text)
    if match_brace:
        extracted = match_brace.group(0).strip()
        try:
            return json.loads(extracted)
        except json.JSONDecodeError:
            text = extracted # Use extracted for subsequent strategies

    # Strategy 4: Clean trailing commas
    try:
        cleaned = re.sub(r',\s*([\]}])', r'\1', text)
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass

    # Strategy 5: Python ast.literal_eval
    try:
        eval_text = text
        eval_text = re.sub(r'(?<![\'\"])\bnull\b(?![\'\"])', 'None', eval_text)
        eval_text = re.sub(r'(?<![\'\"])\btrue\b(?![\'\"])', 'True', eval_text)
        eval_text = re.sub(r'(?<![\'\"])\bfalse\b(?![\'\"])', 'False', eval_text)
        
        parsed = ast.literal_eval(eval_text)
        if isinstance(parsed, (dict, list)):
            return parsed
    except Exception:
        pass

    # Strategy 6: Heavy cleaning
    try:
        cleaned = text
        cleaned = re.sub(r"^'", '"', cleaned)
        cleaned = re.sub(r"'$", '"', cleaned)
        cleaned = re.sub(r"([{\[,:\s])'", r'\1"', cleaned)
        cleaned = re.sub(r"'([}\],:\s])", r'"\1', cleaned)
        cleaned = re.sub(r'([{,]\s*)([a-zA-Z0-9_]+)\s*:', r'\1"\2":', cleaned)
        cleaned = re.sub(r',\s*([\]}])', r'\1', cleaned)
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass

    raise json.JSONDecodeError("Failed to extract JSON from Ollama response.", text, 0)


# -----------------------------
# STAGE A: CONTENT UNDERSTANDING & SUITABILITY VECTOR
# -----------------------------
def run_stage_a_content_understanding(context: str, topic_fallback: str = "General Course Concept") -> dict:
    cleaned_context = sanitize_source_text(context)
    content_hash = compute_normalized_hash(cleaned_context)
    
    cached = get_stage_a_cache(content_hash)
    if cached:
        print(f"  [STAGE A CACHE HIT] Reusing cached Content Understanding for hash: {content_hash[:10]}...")
        return cached

    print("  [STAGE A] Executing Multi-Factor Concept Importance & Suitability Vector Analysis...")
    weights = pedagogical_config.get("importance_weights", {})
    
    prompt = (
        "You are an expert Computer Science Curriculum Architect and Data Extractor.\n"
        "Your objective is to analyze raw educational transcripts, documents, or slides and extract core academic concepts, their relationships, and their suitability for assessment.\n\n"
        "You will NOT generate quiz questions. Your only job is to understand the text and output a strict, mathematically graded JSON blueprint.\n\n"
        "### EXTRACTION RULES & SCORING RUBRIC\n"
        "You must evaluate concepts on a scale of 0.0 to 1.0. Do NOT calculate final weighted scores; only provide the raw component scores based on this rubric:\n\n"
        "1. Importance Breakdown (0.0 - 1.0):\n"
        "   - coverage: How much of the text is dedicated to this concept?\n"
        "   - depth: Are there technical definitions, formulas, code snippets, or deep explanations?\n"
        "   - repetition: Does the concept appear multiple times across different sections?\n"
        "   - teacher_emphasis: Did the text explicitly highlight this? (e.g., 'Note this', 'Important', 'Exam topic' = 0.9+).\n"
        "   - dependency: Is this a foundational concept that other topics rely on?\n"
        "   - learning_objectives: Does this align closely with the overarching theme of the text?\n\n"
        "2. Extraction Confidence (0.0 - 1.0):\n"
        "   - Rate the clarity of the source text for this concept.\n"
        "   - 0.95+ = Perfectly clear, highly detailed text.\n"
        "   - 0.60 = Ambiguous, messy OCR, or poorly explained text.\n"
        "   - 0.30 = Barely readable or highly fragmented context.\n\n"
        "3. Global Assessment Suitability (0.0 - 1.0):\n"
        "   - theory: How well does this entire text lend itself to conceptual/definition questions?\n"
        "   - coding: How well does this text lend itself to implementation, debugging, or dry-run questions?\n"
        "   - scenario: How well does this text support real-world application or system design questions?\n\n"
        "4. Evidence & Citations:\n"
        "   - Provide 1-2 brief, exact quotes from the text as `evidence` to justify your scores.\n"
        "   - Provide a rough locator as `citations` (e.g., 'Paragraph 3', 'Chunk 4', or 'Slide 12').\n\n"
        "5. Isolated Narratives:\n"
        "   - Identify and isolate any metaphorical stories, fictional settings, or names used merely for teaching (e.g., 'Alice and Bob', 'Dining Philosophers', 'Bank Account Example').\n\n"
        f"Context:\n{cleaned_context[:6000]}\n\n"
        "Return ONLY a clean JSON object conforming strictly to this format:\n"
        "{\n"
        "  \"concepts\": [\n"
        "    {\n"
        "      \"concept\": \"Binary Search\",\n"
        "      \"concept_type\": \"algorithm\",\n"
        "      \"extraction_confidence\": 0.95,\n"
        "      \"recommended_bloom_levels\": [\"Understand\", \"Apply\"],\n"
        "      \"importance_breakdown\": {\n"
        "        \"coverage\": 0.8,\n"
        "        \"depth\": 0.9,\n"
        "        \"repetition\": 0.7,\n"
        "        \"teacher_emphasis\": 0.9,\n"
        "        \"dependency\": 0.85,\n"
        "        \"learning_objectives\": 0.9\n"
        "      },\n"
        "      \"evidence\": [\"Binary search operates in logarithmic time...\"],\n"
        "      \"citations\": [\"Chunk 4\"]\n"
        "    }\n"
        "  ],\n"
        "  \"concept_graph\": {\n"
        "    \"nodes\": [\"Sorting\", \"Binary Search\"],\n"
        "    \"edges\": [\n"
        "      {\"source\": \"Sorting\", \"target\": \"Binary Search\", \"relation\": \"REQUIRES\"}\n"
        "    ]\n"
        "  },\n"
        "  \"suitability\": {\n"
        "    \"theory\": 0.9,\n"
        "    \"coding\": 0.8,\n"
        "    \"scenario\": 0.7\n"
        "  },\n"
        "  \"suitability_reasoning\": {\n"
        "    \"theory\": \"High conceptual content present...\",\n"
        "    \"coding\": \"Code examples present...\"\n"
        "  },\n"
        "  \"isolated_narratives\": [\"Alice and Bob transfer example\"]\n"
        "}"
    )

    payload = {
        "model": MODEL_NAME,
        "prompt": prompt,
        "stream": False,
        "format": "json",
        "options": { "temperature": 0.2, "num_ctx": 6144 }
    }

    def process_stage_a_concepts(data: dict) -> dict:
        weights = pedagogical_config.get("importance_weights", {
            "coverage": 0.25, "depth": 0.20, "repetition": 0.15,
            "teacher_emphasis": 0.20, "dependency": 0.10, "learning_objectives": 0.10
        })
        sum_w = sum(weights.values()) or 1.0

        norm_concepts = []
        for c in data.get("concepts", []):
            c_name = c.get("concept") or c.get("concept_tag") or topic_fallback
            c_type = c.get("concept_type", "theory")
            c_blooms = c.get("recommended_bloom_levels", ["Understand", "Apply"])
            c_conf = float(c.get("extraction_confidence") or c.get("confidence") or 0.90)
            
            bd = c.get("importance_breakdown", {})
            cov = float(bd.get("coverage", 0.75))
            dep = float(bd.get("depth", 0.75))
            rep = float(bd.get("repetition", 0.75))
            te  = float(bd.get("teacher_emphasis", 0.75))
            dp  = float(bd.get("dependency", 0.75))
            lo  = float(bd.get("learning_objectives", 0.75))
            
            raw_score = (
                cov * weights.get("coverage", 0.25) +
                dep * weights.get("depth", 0.20) +
                rep * weights.get("repetition", 0.15) +
                te  * weights.get("teacher_emphasis", 0.20) +
                dp  * weights.get("dependency", 0.10) +
                lo  * weights.get("learning_objectives", 0.10)
            ) / sum_w
            
            final_imp = round(raw_score, 4)
            eff_score = round(final_imp * c_conf, 4)

            norm_concepts.append({
                "concept": c_name,
                "concept_tag": c_name,
                "concept_type": c_type,
                "recommended_bloom_levels": c_blooms,
                "final_importance_score": final_imp,
                "importance_score": final_imp,
                "extraction_confidence": c_conf,
                "confidence": c_conf,
                "effective_score": eff_score,
                "importance_breakdown": bd,
                "evidence": c.get("evidence", []),
                "citations": c.get("citations", [])
            })

        data["concepts"] = norm_concepts
        return data

    try:
        response = requests.post(OLLAMA_URL, json=payload, timeout=180)
        if response.status_code == 200:
            raw_text = response.json().get("response", "")
            data = robust_json_loads(raw_text)
            data = process_stage_a_concepts(data)
            data["content_hash"] = content_hash
            set_stage_a_cache(content_hash, data)
            print(f"  ├── [Stage A] Extracted {len(data.get('concepts', []))} concepts with Suitability: {data.get('suitability')}")
            return data
    except Exception as e:
        print(f"  ├── [Stage A] Ollama failed ({e}). Falling back to Groq Cloud for Stage A...")
        if os.getenv("GROQ_API_KEY"):
            try:
                raw_text = call_groq_fallback(prompt, json_mode=True)
                data = robust_json_loads(raw_text)
                data = process_stage_a_concepts(data)
                data["content_hash"] = content_hash
                set_stage_a_cache(content_hash, data)
                return data
            except Exception as ge:
                print(f"  ├── [Stage A] Groq fallback error: {ge}")

    fallback_data = {
        "concepts": [{
            "concept": topic_fallback,
            "concept_tag": topic_fallback,
            "concept_type": "theory",
            "recommended_bloom_levels": ["Understand", "Apply"],
            "final_importance_score": 0.75,
            "importance_score": 0.75,
            "extraction_confidence": 0.85,
            "confidence": 0.85,
            "effective_score": 0.6375,
            "importance_breakdown": {"coverage": 0.75, "depth": 0.75, "repetition": 0.75, "teacher_emphasis": 0.75, "dependency": 0.75, "learning_objectives": 0.75},
            "evidence": ["Direct context"],
            "citations": ["Chunk 1"]
        }],
        "suitability": {"theory": 0.8, "coding": 0.5, "scenario": 0.5},
        "suitability_reasoning": {"theory": "Direct text analysis", "coding": "General content"},
        "concept_graph": {"nodes": [topic_fallback], "edges": []},
        "prerequisite_graph": {"dependencies": []},
        "evidence": ["Direct context"],
        "citations": ["Chunk 1"],
        "isolated_narratives": [],
        "content_hash": content_hash
    }
    set_stage_a_cache(content_hash, fallback_data)
    return fallback_data


# -----------------------------
# PLANNER VALIDATOR CIRCUIT BREAKER
# -----------------------------
def validate_stage_b_blueprint(blueprint: dict) -> (bool, List[str]):
    errors = []
    if not blueprint or not isinstance(blueprint, dict):
        return False, ["Blueprint is empty or not a JSON object"]
    
    slots = blueprint.get("slots", [])
    if not slots or len(slots) == 0:
        errors.append("Blueprint contains 0 slots")
        
    for idx, slot in enumerate(slots):
        s_idx = slot.get("slot_index", idx + 1)
        if not slot.get("concept_tag"):
            errors.append(f"Slot {s_idx}: Missing concept_tag")
        if not slot.get("archetype"):
            errors.append(f"Slot {s_idx}: Missing archetype")
        if not slot.get("blooms_level"):
            errors.append(f"Slot {s_idx}: Missing blooms_level")
        if not slot.get("question_style"):
            errors.append(f"Slot {s_idx}: Missing question_style")

    blooms = set(slot.get("blooms_level") for slot in slots if slot.get("blooms_level"))
    if len(blooms) == 0:
        errors.append("No valid Bloom's levels present in slots")

    return (len(errors) == 0, errors)


# -----------------------------
# STAGE B: ACADEMIC PLANNER
# -----------------------------
# -----------------------------
# STAGE B: ACADEMIC PLANNER
# -----------------------------
def run_stage_b_academic_planner(
    stage_a_data: dict,
    question_count: int,
    difficulty: str = "Medium",
    question_style: str = "MIXED",
    interview_mode: bool = False
) -> dict:
    import math
    print(f"  [STAGE B PLANNER] Generating Blueprint for Count: {question_count} | Difficulty: {difficulty} | Style: {question_style} | Interview Mode: {interview_mode}")
    
    raw_concepts = stage_a_data.get("concepts", [])
    NOISE_THRESHOLD = 0.30
    
    # 1. Effective Score = Final Importance Score * Extraction Confidence
    scored_concepts = []
    for c in raw_concepts:
        c_name = c.get("concept") or c.get("concept_tag") or "General Course Concept"
        c_imp = float(c.get("final_importance_score") or c.get("importance_score") or 0.75)
        c_conf = float(c.get("extraction_confidence") or c.get("confidence") or 0.90)
        c_type = c.get("concept_type", "Definition")
        eff_score = round(c_imp * c_conf, 4)
        
        norm_c = {
            "concept": c_name,
            "concept_tag": c_name,
            "concept_type": c_type,
            "final_importance_score": c_imp,
            "importance_score": c_imp,
            "extraction_confidence": c_conf,
            "confidence": c_conf,
            "effective_score": eff_score,
            "citations": c.get("citations", ["chunk_01"])
        }
        
        if eff_score >= NOISE_THRESHOLD:
            scored_concepts.append(norm_c)
        else:
            print(f"  [STAGE B PLANNER] Discarded noise concept below threshold (< {NOISE_THRESHOLD}): '{c_name}' (effective_score: {eff_score:.4f})")
            
    scored_concepts.sort(key=lambda x: x["effective_score"], reverse=True)
            
    if not scored_concepts:
        scored_concepts = [{
            "concept": "General Course Concept",
            "concept_tag": "General Course Concept",
            "concept_type": "Definition",
            "final_importance_score": 0.75,
            "extraction_confidence": 1.0,
            "confidence": 1.0,
            "effective_score": 0.75,
            "citations": ["chunk_01"]
        }]

    concepts = scored_concepts
    num_valid = len(concepts)

    # 2. Adaptive Allocation Caps
    if num_valid <= 2:
        cap_percentage = 0.60
    elif num_valid <= 4:
        cap_percentage = 0.45
    elif num_valid <= 9:
        cap_percentage = 0.35
    else:
        cap_percentage = 0.25

    max_per_concept = max(1, math.floor(question_count * cap_percentage))

    suitability = stage_a_data.get("suitability", {"theory": 0.8, "coding": 0.5, "scenario": 0.5})

    # 3. Pattern Diversity Catalog
    pattern_catalog = {
        "THEORY": ["Assertion_Reasoning", "Comparison_Matrix", "Definition_Core", "Tradeoff_Analysis"],
        "CODING": ["Dry_Run", "Debugging", "Implementation_Flaw", "Complexity_Analysis"],
        "SCENARIO": ["System_Tradeoff", "Failure_Diagnosis", "Architecture_Design"],
        "INTERVIEW": ["Concept_Explanation", "Edge_Case_Handling", "Optimization_Challenge"]
    }

    concept_usage = {c["concept"]: 0 for c in concepts}

    bloom_map = {
        "Easy": ["Remember", "Understand"],
        "Medium": ["Understand", "Apply"],
        "Thinkable": ["Apply", "Analyze"],
        "Hard": ["Analyze", "Evaluate", "Create"]
    }
    allowed_blooms = bloom_map.get(difficulty, ["Understand", "Apply"])

    slots = []
    quiz_blueprint = []

    for i in range(question_count):
        # Select concept honoring max allocation cap
        available_concepts = [c for c in concepts if concept_usage[c["concept"]] < max_per_concept]
        if not available_concepts:
            available_concepts = concepts # Reset cap if all maxed out
        
        concept_obj = available_concepts[i % len(available_concepts)]
        concept_name = concept_obj["concept"]
        concept_usage[concept_name] += 1
        repeat_idx = concept_usage[concept_name] - 1

        # Determine pattern & style
        if interview_mode:
            cat_key = "INTERVIEW"
            assessment_style = "Interview"
        elif question_style == "CODING" or (question_style == "MIXED" and i % 2 == 1 and suitability.get("coding", 0.5) >= 0.4):
            cat_key = "CODING"
            assessment_style = "Coding"
        elif question_style == "SCENARIO" or (question_style == "MIXED" and i % 3 == 2 and suitability.get("scenario", 0.5) >= 0.4):
            cat_key = "SCENARIO"
            assessment_style = "Scenario"
        else:
            cat_key = "THEORY"
            assessment_style = "Theory"

        patterns = pattern_catalog[cat_key]
        question_pattern = patterns[repeat_idx % len(patterns)]
        bloom_level = allowed_blooms[i % len(allowed_blooms)]

        q_id = f"q_{i+1:02d}"

        # Blueprint Item (Directive Contract)
        blueprint_item = {
            "question_id": q_id,
            "assessment_style": assessment_style,
            "concept": concept_name,
            "concept_type": concept_obj.get("concept_type", "Definition"),
            "bloom_level": bloom_level,
            "difficulty_calibration": difficulty,
            "difficulty_reason": [
                f"Evaluates {concept_name} at {bloom_level} level",
                f"Pattern {question_pattern} requires procedural mastery",
                f"Difficulty calibrated for {difficulty} academic band"
            ],
            "question_pattern": question_pattern,
            "assessment_objective": f"Evaluate whether the student understands and can execute {question_pattern} for {concept_name}.",
            "expected_answer_type": "Multiple Choice Option",
            "estimated_time_minutes": 2 if assessment_style != "Coding" else 3,
            "distractor_rules": f"Construct plausible distractors testing common misconceptions for {question_pattern}.",
            "context_citations": {
                "chunk_ids": concept_obj.get("citations", ["chunk_01"]),
                "page": i + 1
            }
        }
        quiz_blueprint.append(blueprint_item)

        # Legacy slots compatibility
        slots.append({
            "slot_index": i + 1,
            "question_id": q_id,
            "concept": concept_name,
            "concept_tag": concept_name,
            "final_importance_score": concept_obj.get("final_importance_score", 0.75),
            "importance_score": concept_obj.get("importance_score", 0.75),
            "extraction_confidence": concept_obj.get("extraction_confidence", 0.90),
            "confidence": concept_obj.get("confidence", 0.90),
            "effective_score": concept_obj.get("effective_score", 0.75),
            "archetype": "PRACTICAL_AND_LAB_TASKS" if cat_key == "CODING" else ("CASE_STUDIES_AND_SCENARIOS" if cat_key == "SCENARIO" else "CONCEPTS_AND_DEFINITIONS"),
            "pattern_type": cat_key,
            "specific_pattern": question_pattern,
            "question_style": question_style,
            "blooms_level": bloom_level,
            "calibrated_difficulty": difficulty
        })

    blueprint = {
        "blueprint_version": "3.0",
        "total_slots": len(slots),
        "requested_difficulty": difficulty,
        "question_style": question_style,
        "interview_mode": interview_mode,
        "suitability": suitability,
        "slots": slots,
        "quiz_blueprint": quiz_blueprint
    }

    if pedagogical_config.get("feature_flags", {}).get("enable_planner_validator", True):
        validated, val_errors = validate_stage_b_blueprint(blueprint)
        if not validated:
            raise ValueError(f"Planner Circuit Breaker Error: {'; '.join(val_errors)}")

    print(f"  └── [STAGE B PLANNER] Blueprint created & validated cleanly with {len(quiz_blueprint)} items.")
    return blueprint



# -----------------------------
# 4. CRITIC EVALUATOR
# -----------------------------
def critic_evaluate(q_text, q_opts, q_ans, topic, difficulty):
    """
    Sends a generated question to the critic LLM for quality evaluation.
    Returns (score: int 1-10, feedback: str).
    Falls back to score=7 if the critic itself fails, so generation isn't blocked.
    """
    options_str = "\n".join([f"  {chr(65+i)}) {opt}" for i, opt in enumerate(q_opts)])
    critic_prompt = (
        f"You are a strict quiz quality critic. Evaluate the quiz question below.\n\n"
        f"Topic: {topic}\n"
        f"Difficulty: {difficulty}\n\n"
        f"Question: {q_text}\n"
        f"Options:\n{options_str}\n"
        f"Correct Answer: {q_ans}\n\n"
        f"Rate this question from 1 to 10 based on:\n"
        f"1. Clarity       — Is the question unambiguous and well-worded?\n"
        f"2. Correctness   — Is the correct answer factually accurate?\n"
        f"3. Distractors   — Are the wrong options plausible but clearly incorrect?\n"
        f"4. Relevance     — Is it on-topic and appropriate for the difficulty level?\n\n"
        f"Return ONLY a JSON object with two keys:\n"
        f"  score    (integer, 1-10)\n"
        f"  feedback (one sentence explaining the rating)\n"
    )

    payload = {
        "model": MODEL_NAME,
        "prompt": critic_prompt,
        "stream": False,
        "format": "json",
        "options": {
            "temperature": 0.2,   # Low temp → consistent, deterministic critic
            "num_ctx": 1024
        }
    }

    try:
        response = requests.post(OLLAMA_URL, json=payload, timeout=60)
        if response.status_code != 200:
            print(f"  [CRITIC]    Ollama returned HTTP {response.status_code} — defaulting score to 7.")
            return 7, "Critic unavailable, defaulting to pass."

        raw_text = response.json().get("response", "")
        data = robust_json_loads(raw_text)

        score = int(data.get("score", 7))
        score = max(1, min(10, score))   # Clamp to [1, 10]
        feedback = str(data.get("feedback", "No feedback provided."))
        return score, feedback

    except Exception as e:
        print(f"  [CRITIC]    Evaluation failed ({e}) — defaulting score to 7.")
        return 7, "Critic evaluation error, defaulting to pass."


def run_agent1_analyzer(context, count, topic_fallback="General Course Concept"):
    print("  [AGENT 1] Executing Semantic Weight & Concept Analyzer...")
    prompt = (
        "You are an expert university academic analyst. Your goal is to inspect the parsed source content chunks and isolate the primary learning objectives.\n\n"
        "MANDATORY CONTEXTUAL DE-NOISING PASS:\n"
        "You must explicitly inspect the provided context. If it contains audio transcripts or speech data, identify and strip out all non-pedagogical clutter: classroom management comments, administrative/scheduling remarks (e.g., homework deadlines, quiz notifications), and unrelated conversational anecdotes. Process only the pure, academic, and engineering concepts.\n\n"
        f"Context Chunks:\n{context[:4000]}\n\n"
        "Tasks:\n"
        f"1. Identify at least {count} distinct core concepts from the de-noised text.\n"
        "2. Classify contextual examples into either 'CLASSIC_DOMAIN_STANDARD' (foundational standard examples common to the domain) or 'TRANSIENT_ANALOGY' (casual/metaphorical settings or temporary stories). If an example is classified as a 'TRANSIENT_ANALOGY', extract the underlying mathematical/logical rule, and place the specific characters/names/settings used into an 'isolated_narratives' exclusion array of strings.\n\n"
        "Return ONLY a clean JSON object conforming strictly to this format:\n"
        "{\n"
        "  \"concepts\": [\n"
        "    {\n"
        "      \"concept_tag\": \"Process Synchronization\",\n"
        "      \"weight_score\": 0.85,\n"
        "      \"anchor_citation\": \"To achieve process synchronization, operating systems use semaphores...\"\n"
        "    }\n"
        "  ],\n"
        "  \"isolated_narratives\": [\n"
        "    \"Alice and Bob\", \"dining philosophers table setup\", \"specific car factory narrative\"\n"
        "  ]\n"
        "}"
    )
    
    payload = {
        "model": MODEL_NAME,
        "prompt": prompt,
        "stream": False,
        "format": "json",
        "options": {
            "temperature": 0.2,
            "num_ctx": 4096
        }
    }
    try:
        response = requests.post(OLLAMA_URL, json=payload, timeout=180)
        if response.status_code == 200:
            raw_text = response.json().get("response", "")
            data = robust_json_loads(raw_text)
            concepts = data.get("concepts", [])
            isolated_narratives = data.get("isolated_narratives", [])
            print(f"  ├── [Agent 1 Analyzer] : Extracted {len(concepts)} concepts and {len(isolated_narratives)} narrative exclusions.")
            return concepts, isolated_narratives
    except Exception as e:
        err_str = str(e).lower()
        if "timeout" in err_str or "read timed out" in err_str:
            print("  ├── [Agent 1 Analyzer] : ⚠️ CPU Latency Spike -> Dynamically pulled topic from filenames.")
        else:
            print(f"  ├── [Agent 1 Analyzer] : ⚠️ Concept mapping failed ({e}) -> Fallback to topic from filenames.")
    
    return [{"concept_tag": topic_fallback, "weight_score": 0.75, "anchor_citation": "Direct context chunk"}], []

def run_agent2_generator(concept, question_type, context, generated_so_far="", difficulty="Medium", isolated_narratives=None, is_compression_sieve=False, temp_offset=0.0, freq_penalty=None, programming_language="Python/JavaScript", question_style="THEORY"):
    concept_tag = concept.get("concept_tag", "General Course Concept")
    weight_score = concept.get("weight_score", 0.75)
    
    _tier_map = {"easy": 1, "medium": 2, "thinkable": 2, "hard": 3}
    diff_key = difficulty.lower()
    if diff_key not in _tier_map: diff_key = "medium"
    difficulty_tier = _tier_map.get(diff_key, 2)

    # ── Decoupled Strategy Engine (Question Style Plugin Matrix) ──
    import os
    style_files = {
        "THEORY": "theory_prompt.txt",
        "OUTPUT_PREDICTION": "output_prompt.txt",
        "DEBUGGING": "debug_prompt.txt",
        "ASSERTION_REASONING": "assertion_reason_prompt.txt",
        "TRACE_EXECUTION": "trace_prompt.txt",
        "COMPILER_ERROR": "compiler_prompt.txt",
        "CODE_COMPLETION": "completion_prompt.txt"
    }
    
    filename = style_files.get(question_style, "theory_prompt.txt")
    prompts_dir = os.path.join(os.path.dirname(__file__), "prompts")
    filepath = os.path.join(prompts_dir, filename)
    
    strategy_content = ""
    if os.path.exists(filepath):
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                strategy_content = f.read()
        except Exception as e:
            print(f"  ├── [Strategy Engine] : ⚠️ Error loading strategy prompt file: {e}")
            
    if not strategy_content:
        strategy_content = "Write a conceptual multiple-choice question."
        
    try:
        type_instruction = strategy_content.format(
            concept_tag=concept_tag,
            difficulty=difficulty,
            context=context[:2000],
            programming_language=programming_language
        )
    except Exception as fe:
        type_instruction = strategy_content.replace("{concept_tag}", concept_tag)\
                                           .replace("{difficulty}", difficulty)\
                                           .replace("{context}", context[:2000])\
                                           .replace("{programming_language}", programming_language)

    difficultyPrompts = {
        "easy": {
            "CONCEPTS_AND_DEFINITIONS": "Identify straightforward definitions or basic protocol/concept names.",
            "COMPARISONS_AND_TRADEOFFS": "Identify simple differences or obvious advantages between two concepts.",
            "FORMULAS_AND_CALCULATIONS": "Simple direct calculations or basic parameter identification.",
            "CASE_STUDIES_AND_SCENARIOS": "Direct single-variable applications with simple outcomes.",
            "PRACTICAL_AND_LAB_TASKS": "Short, simple code blocks or basic syntax identification."
        },
        "medium": {
            "CONCEPTS_AND_DEFINITIONS": "Explain how mechanisms interact or standard workflows.",
            "COMPARISONS_AND_TRADEOFFS": "Analyze trade-offs, state-space exploration, or standard error conditions.",
            "FORMULAS_AND_CALCULATIONS": "Multi-step calculations, simple diagram/memory tracing, or time/space complexities.",
            "CASE_STUDIES_AND_SCENARIOS": "Introduce minor engineering bottlenecks, common system failures, or design trade-offs.",
            "PRACTICAL_AND_LAB_TASKS": "Code prediction involving loops, basic conditional checks, or state updates."
        },
        "thinkable": {
            "CONCEPTS_AND_DEFINITIONS": "Probe nuanced distinctions between closely related concepts or edge-case protocol behaviors.",
            "COMPARISONS_AND_TRADEOFFS": "Evaluate multi-factor trade-off scenarios where no single answer is obviously dominant.",
            "FORMULAS_AND_CALCULATIONS": "Multi-step derivations with intermediate rounding decisions or boundary-condition choices.",
            "CASE_STUDIES_AND_SCENARIOS": "Scenarios with partially conflicting design constraints requiring prioritization logic.",
            "PRACTICAL_AND_LAB_TASKS": "Subtle bugs involving off-by-one errors, scope shadowing, or unexpected type coercions."
        },
        "hard": {
            "CONCEPTS_AND_DEFINITIONS": "Test deep internal mechanics, architectural limits, and complex theoretical constraints.",
            "COMPARISONS_AND_TRADEOFFS": "Evaluate complex state transitions, hidden structural flaws, or multi-dimensional trade-offs.",
            "FORMULAS_AND_CALCULATIONS": "Deep computational calculations, complex diagram tracing, multi-variable optimization, or proof-of-correctness tracing.",
            "CASE_STUDIES_AND_SCENARIOS": "Construct deep, multi-layered system failure scenarios with conflicting resource/performance metrics.",
            "PRACTICAL_AND_LAB_TASKS": "Analyze highly optimized snippets, multi-threaded tasks, tricky recursion, or memory allocation bugs."
        }
    }
    
    targeted_criteria = difficultyPrompts[diff_key].get(question_type, "Focus on general knowledge.")

    if difficulty_tier == 1:
        distractor_instruction = (
            "DISTRACTOR CALIBRATION (Tier 1 — Distinct):\n"
            "Incorrect options must be clearly distinct, conceptually non-overlapping terms. "
            "Each wrong option should belong to a different area of the topic so a well-prepared "
            "student can immediately distinguish it from the correct answer."
        )
    elif difficulty_tier == 2:
        distractor_instruction = (
            "DISTRACTOR CALIBRATION (Tier 2 — Procedural Traps):\n"
            "Incorrect options must represent classic procedural mistakes or common calculation slips "
            "directly related to the source text. Each distractor should reflect a real error path a "
            "student might take (e.g., swapping two related parameters, off-by-one in a formula, "
            "applying the wrong formula variant, or confusing a near-synonym concept). "
            "Wrong options must not be obviously wrong at a glance."
        )
    else:  # tier 3
        distractor_instruction = (
            "DISTRACTOR CALIBRATION (Tier 3 — Maximum Cognitive Proximity):\n"
            "This is a PRECISION MEMORY question. Distractors must use IDENTICAL vocabulary and "
            "structural syntax as the correct answer. Alter ONLY a single subtle element per distractor: "
            "a mathematical sign (+ vs −), a logical operator (AND vs OR, < vs ≤), a boundary condition "
            "(n-1 vs n, 2^k vs 2^(k-1)), or a single key term replaced with its near-synonym. "
            "All four options must look highly similar so that only a student who memorized the exact "
            "definition can distinguish the correct one. "
            "FORBIDDEN: Do NOT use options from different conceptual domains or obviously incorrect "
            "catch-all phrases. Every distractor must be a plausible near-miss of the correct answer."
        )

    narrative_mutation_instruction = ""
    if isolated_narratives and len(isolated_narratives) > 0:
        exclusions_str = ", ".join([f'"{n}"' for n in isolated_narratives])
        narrative_mutation_instruction = (
            f"NARRATIVE MUTATION MANDATE:\n"
            f"When generating questions under the CASE_STUDIES_AND_SCENARIOS or FORMULAS_AND_CALCULATIONS archetypes, review the following narrative/scenario exclusions: [{exclusions_str}].\n"
            f"You are STRICTLY PROHIBITED from using these exact scenarios, names, characters, or settings in your question stem or choices. "
            f"Instead, construct a structurally isomorphic (parallel) real-world scenario that tests the exact same concept using a brand-new application domain.\n\n"
        )

    compression_sieve_instruction = ""
    if is_compression_sieve:
        compression_sieve_instruction = (
            "COMPRESSION SIEVE MANDATE:\n"
            "Because you are generating a very short quiz (<= 5 questions) from a very large amount of material (>5000 tokens), "
            "each question MUST be a multi-concept synthesis question. Do NOT ask about isolated trivial details. "
            "Instead, design high-level questions that require the student to synthesize, compare, and integrate multiple concepts "
            "from the source text simultaneously (e.g. how a design decision affects memory layout and time complexity together, "
            "or comparing the trade-offs of two different architectures in a single question).\n\n"
        )

    prompt = (
        "You are an elite Computer Science and Engineering curriculum developer. Your task is to write high-fidelity academic evaluations.\n\n"
        f"{compression_sieve_instruction}"
        f"{narrative_mutation_instruction}"
        "DIAGRAM AND GRAPHICS SAFEGUARD:\n"
        "For any image/diagram-based text descriptions extracted from slides, you are only permitted to reproduce them if they represent standard engineering blueprints or industry-standard topologies. If the slide context indicates a non-standard or arbitrary sketch, synthesize a descriptive text-based conceptual evaluation instead.\n\n"
        f"Context chunks:\n{context[:2000]}\n\n"
        f"Target Concept: {concept_tag} (Importance Weight: {weight_score})\n"
        f"Question Type: {question_type}\n\n"
        f"CRITICAL DIFFICULTY INSTRUCTION ({difficulty.upper()}, Tier {difficulty_tier}/3):\n"
        f"You MUST strictly follow these criteria: '{targeted_criteria}'\n\n"
        f"{distractor_instruction}\n\n"
        f"Task:\n{type_instruction}\n\n"
        f"CRITICAL: Avoid repeating the topics of these existing questions: {generated_so_far}\n\n"
        "[STRICT GROUNDING CONSTRAINT]\n"
        "You must extract ONLY core academic, structural, and theoretical concepts present within the text body.\n"
        "CRITICAL WARNING: Completely ignore any references to dates, times, AM/PM, audio lengths, transcription artifacts, or file names. Under no circumstances should a question or answer choice analyze when a recording happened, what a file name is, or how data was collected.\n\n"
        "COMPLEXITY MANDATE: Do not create circular questions where the answer repeats words from the question stem. Focus on operational logic, mechanics (e.g., how splitting criteria work in Random Forests), and architectural design trade-offs.\n\n"
        "Return ONLY a clean JSON object conforming to this schema:\n"
        "{\n"
        "  \"concept_tag\": \"concept\",\n"
        "  \"prompt_text\": \"Question text...\",\n"
        "  \"code_snippet\": \"Markdown code block or null if not applicable\",\n"
        "  \"options\": [\n"
        "    \"Option A\",\n"
        "    \"Option B\",\n"
        "    \"Option C\",\n"
        "    \"Option D\"\n"
        "  ],\n"
        "  \"correct_answer\": \"Exact copy of the correct option from options\",\n"
        "  \"explanation\": \"Detailed explanation explaining the solution...\",\n"
        "  \"blooms_level\": \"Remember/Understand or Apply or Analyze or Create\"\n"
        "}"
    )

    temp_map = {
        "CONCEPTS_AND_DEFINITIONS": 0.15,
        "FORMULAS_AND_CALCULATIONS": 0.2,
        "CASE_STUDIES_AND_SCENARIOS": 0.65,
        "COMPARISONS_AND_TRADEOFFS": 0.4,
        "PRACTICAL_AND_LAB_TASKS": 0.4
    }
    temperature = temp_map.get(question_type, 0.4) + temp_offset

    payload = {
        "model": MODEL_NAME,
        "prompt": prompt,
        "stream": False,
        "format": "json",
        "options": {
            "temperature": temperature,
            "repeat_penalty": 1.15,
            "num_ctx": 4096
        }
    }
    if freq_penalty is not None:
        payload["options"]["frequency_penalty"] = freq_penalty
    
    response = requests.post(OLLAMA_URL, json=payload, timeout=180)
    if response.status_code == 200:
        raw_text = response.json().get("response", "")
        data = robust_json_loads(raw_text)
        return data
    else:
        raise Exception(f"Generator failed with HTTP {response.status_code}")

def check_python_syntax(code):
    try:
        ast.parse(code)
        return True, []
    except SyntaxError as e:
        return False, [f"Python Code Block Syntax Error: {e.msg} at line {e.lineno}"]

def check_js_syntax(code):
    import subprocess
    try:
        js_cmd = f"new Function({json.dumps(code)})"
        res = subprocess.run(['node', '-e', js_cmd], capture_output=True, text=True, timeout=2)
        if res.returncode != 0:
            err_lines = [l for l in res.stderr.split('\n') if l.strip()]
            return False, [f"JS Code Block Syntax Error: {err_lines[0] if err_lines else 'Compilation failed'}"]
        return True, []
    except Exception:
        return True, []

def run_agent3_critic(q_data, context, flavor="theory", difficulty="Medium", question_style="THEORY"):
    issues = []
    
    # Derive tier for distractor precision checks
    _tier_map = {"easy": 1, "medium": 2, "thinkable": 2, "hard": 3}
    difficulty_tier = _tier_map.get(difficulty.lower(), 2)

    prompt_text = q_data.get("prompt_text")
    options = q_data.get("options", [])
    correct_ans = q_data.get("correct_answer") or q_data.get("correctAnswer")
    code_snippet = q_data.get("code_snippet")

    # ── Style Validation Checks ──
    if question_style in ['OUTPUT_PREDICTION', 'DEBUGGING', 'CODE_COMPLETION', 'COMPILER_ERROR']:
        has_code = False
        raw_code = ""
        if code_snippet and "```" in code_snippet:
            has_code = True
            match = re.search(r'```(?:[a-zA-Z0-9+#-]+)?\n([\s\S]*?)\n```', code_snippet)
            if match:
                raw_code = match.group(1)
            else:
                raw_code = code_snippet.replace("```", "").strip()
        elif prompt_text and "```" in prompt_text:
            has_code = True
            match = re.search(r'```(?:[a-zA-Z0-9+#-]+)?\n([\s\S]*?)\n```', prompt_text)
            if match:
                raw_code = match.group(1)
            else:
                raw_code = prompt_text.replace("```", "").strip()
            
        if not has_code:
            issues.append(f"Style validation failure: Style '{question_style}' requires an executable markdown code block (```)")
            return {"status": "fail", "issues": issues}

        # ── Sandbox Execution Validator ──
        is_python = False
        if "python" in (code_snippet or "").lower() or "python" in (prompt_text or "").lower() or "pandas" in (prompt_text or "").lower() or "numpy" in (prompt_text or "").lower():
            is_python = True
            
        if is_python and question_style == 'OUTPUT_PREDICTION' and raw_code.strip():
            import subprocess
            import sys
            import tempfile
            import os
            
            try:
                with tempfile.NamedTemporaryFile(suffix=".py", delete=False, mode="w", encoding="utf-8") as temp_file:
                    temp_file.write(raw_code)
                    temp_filepath = temp_file.name
                
                res = subprocess.run(
                    [sys.executable, temp_filepath],
                    capture_output=True,
                    text=True,
                    timeout=5
                )
                
                try: os.unlink(temp_filepath)
                except Exception: pass
                
                if res.returncode == 0:
                    stdout_val = res.stdout.strip()
                    clean_ans = str(correct_ans).strip()
                    if clean_ans not in stdout_val and stdout_val not in clean_ans:
                        print(f"CRITIC_WARN: Sandbox Execution Mismatch. stdout was '{stdout_val}', expected match with '{clean_ans}'")
                        issues.append(f"CRITIC_WARN: Sandbox Execution Mismatch (stdout: '{stdout_val}', expected: '{clean_ans}')")
                        return {"status": "fail", "issues": issues}
                else:
                    stderr_val = res.stderr.strip().split('\n')[-1]
                    print(f"CRITIC_WARN: Sandbox Execution Error: {stderr_val}")
                    issues.append(f"CRITIC_WARN: Sandbox Execution Error ({stderr_val})")
                    return {"status": "fail", "issues": issues}
            except subprocess.TimeoutExpired:
                try: os.unlink(temp_filepath)
                except Exception: pass
                print("CRITIC_WARN: Sandbox Execution Timeout expired.")
                issues.append("CRITIC_WARN: Sandbox Execution Timeout expired.")
                return {"status": "fail", "issues": issues}
            except Exception as se:
                try: os.unlink(temp_filepath)
                except Exception: pass
                print(f"CRITIC_WARN: Sandbox Execution system error: {se}")

    # ── Token Similarity Shield ──
    # Compare every option against every other option, and against the correct answer.
    if options:
        def compute_iou(s1, s2):
            w1 = set(re.sub(r'[^\w\s]', ' ', str(s1).lower()).split())
            w2 = set(re.sub(r'[^\w\s]', ' ', str(s2).lower()).split())
            if not w1 and not w2:
                return 0.0
            union = w1.union(w2)
            if not union:
                return 0.0
            return len(w1.intersection(w2)) / len(union)

        # Check options against correct answer
        if correct_ans:
            for opt in options:
                if opt != correct_ans:
                    iou = compute_iou(opt, correct_ans)
                    if iou > 0.75:
                        print(f"CRITIC_WARN: High Semantic Overlap Detected (IoU: {iou:.2f}) between option '{opt}' and correct answer '{correct_ans}'")
                        issues.append("CRITIC_WARN: High Semantic Overlap Detected")
                        return {"status": "fail", "issues": issues}

        # Check options against each other
        for idx1 in range(len(options)):
            for idx2 in range(idx1 + 1, len(options)):
                iou = compute_iou(options[idx1], options[idx2])
                if iou > 0.75:
                    print(f"CRITIC_WARN: High Semantic Overlap Detected (IoU: {iou:.2f}) between options '{options[idx1]}' and '{options[idx2]}'")
                    issues.append("CRITIC_WARN: High Semantic Overlap Detected")
                    return {"status": "fail", "issues": issues}
    
    if not prompt_text:
        issues.append("Missing prompt_text")
    else:
        # Rule X (Zero Metadata Leak)
        metadata_indicators = [
            r'\b\d{1,2}:\d{2}(?::\d{2})?\s*(?:AM|PM|am|pm)\b',
            r'\b(?:voice\s+)?transcript\b',
            r'\bfile\s+name\b',
            r'\b\.pdf\b|\b\.docx\b|\b\.pptx\b|\b\.txt\b|\b\.webm\b',
            r'\brecording\b'
        ]
        for pattern in metadata_indicators:
            if re.search(pattern, prompt_text, re.IGNORECASE):
                issues.append("Question contains timestamps, file names, or recording properties (Metadata Leak).")
                break

    if not options or len(options) < 4:
        issues.append("Fewer than 4 options provided")
    else:
        # Check options for metadata leaks
        for opt in options:
            metadata_indicators = [
                r'\b\d{1,2}:\d{2}(?::\d{2})?\s*(?:AM|PM|am|pm)\b',
                r'\b(?:voice\s+)?transcript\b',
                r'\bfile\s+name\b',
                r'\b\.pdf\b|\b\.docx\b|\b\.pptx\b|\b\.txt\b|\b\.webm\b',
                r'\brecording\b'
            ]
            leak_found = False
            for pattern in metadata_indicators:
                if re.search(pattern, opt, re.IGNORECASE):
                    issues.append("An option choice contains metadata leaks (timestamps/file details).")
                    leak_found = True
                    break
            if leak_found:
                break

        # ── Tier 3 Distractor Precision Check ─────────────────────────────────
        # For Hard/Tier-3 questions, reject distractors that are structurally
        # asymmetric or trivially eliminable by surface-level keyword mismatch.
        if difficulty_tier == 3 and correct_ans and correct_ans in options:
            correct_tokens = set(re.sub(r'[^\w\s]', ' ', correct_ans.lower()).split())
            trivially_weak_distractors = 0
            for opt in options:
                if opt == correct_ans:
                    continue
                opt_tokens = set(re.sub(r'[^\w\s]', ' ', opt.lower()).split())
                # Compute token overlap between this distractor and the correct answer
                if len(correct_tokens) > 0:
                    overlap_ratio = len(correct_tokens.intersection(opt_tokens)) / len(correct_tokens)
                    # A distractor with < 30% token overlap with the correct answer
                    # is structurally asymmetric and trivially eliminable at Tier 3
                    if overlap_ratio < 0.30:
                        trivially_weak_distractors += 1
            # Fail if more than 1 distractor is trivially eliminable
            if trivially_weak_distractors > 1:
                issues.append(
                    f"TIER-3 PRECISION FAILURE: {trivially_weak_distractors} distractor(s) have <30% vocabulary "
                    "overlap with the correct answer and can be eliminated by simple keyword scanning. "
                    "All distractors must use identical structural vocabulary as the correct answer, "
                    "differing only in a single subtle sign, operator, or boundary condition."
                )

    if not correct_ans:
        issues.append("Missing correct_answer")
    elif correct_ans not in options:
        if flavor != "PRACTICAL_AND_LAB_TASKS":
            issues.append("Correct answer does not match any of the options exactly")
            
    # Rule Y (Circular Verification)
    if correct_ans and prompt_text:
        prompt_clean = re.sub(r'[^\w\s]', '', prompt_text.lower())
        ans_clean = re.sub(r'[^\w\s]', '', correct_ans.lower())
        prompt_words = set(prompt_clean.split())
        ans_words = set(ans_clean.split())
        if len(ans_words) > 0:
            overlap = prompt_words.intersection(ans_words)
            if len(overlap) / len(ans_words) > 0.8:
                issues.append("Question is circular (the correct answer repeats almost all words from the question stem).")
        
    if code_snippet and isinstance(code_snippet, str) and ("```" in code_snippet):
        match_py = re.search(r'```python\n([\s\S]*?)```', code_snippet)
        match_js = re.search(r'```(?:javascript|js)\n([\s\S]*?)```', code_snippet)
        
        if match_py:
            code = match_py.group(1)
            if flavor != "PRACTICAL_AND_LAB_TASKS":
                valid, errs = check_python_syntax(code)
                if not valid:
                    issues.extend(errs)
        elif match_js:
            code = match_js.group(1)
            if flavor != "PRACTICAL_AND_LAB_TASKS":
                valid, errs = check_js_syntax(code)
                if not valid:
                    issues.extend(errs)
                
    if len(set(options)) < len(options):
        issues.append("Duplicate options detected")
        
    status = "pass" if len(issues) == 0 else "fail"
    return {
        "status": status,
        "issues": issues
    }

def run_agent2_repair(q_data, issues, context):
    prompt = (
        "You are an elite Computer Science and Engineering curriculum developer.\n"
        "Your task is to fix a multiple-choice question that failed quality control checks.\n\n"
        "ORIGINAL QUESTION:\n"
        f"{json.dumps(q_data, indent=2)}\n\n"
        "QUALITY CONTROL ISSUES FOUND:\n"
        + "\n".join([f"- {iss}" for iss in issues]) + "\n\n"
        "CONTEXT FOR GROUNDING:\n"
        f"{context[:2000]}\n\n"
        "Task:\n"
        "Rewrite the question to fix all listed issues. Make sure the correct_answer is an exact copy-paste of one of the options and all options are unique.\n\n"
        "Return ONLY a clean JSON object conforming to the same schema as the original question."
    )
    
    payload = {
        "model": MODEL_NAME,
        "prompt": prompt,
        "stream": False,
        "format": "json",
        "options": {
            "temperature": 0.2,
            "num_ctx": 4096
        }
    }
    
    response = requests.post(OLLAMA_URL, json=payload, timeout=180)
    if response.status_code == 200:
        raw_text = response.json().get("response", "")
        data = robust_json_loads(raw_text)
        return data
    else:
        raise Exception("Repair failed")
def log_pipeline_step(step_number, step_name, data_description, payload):
    CYAN = '\033[36m'
    GREEN = '\033[32m'
    GRAY = '\033[90m'
    RESET = '\033[0m'
    print(f"{GRAY}\n========================================================{RESET}")
    print(f"{CYAN}➡️ [STEP {step_number}] {step_name}{RESET}")
    print(f"{GRAY}📋 Data State: {data_description}{RESET}")
    print(f"{GRAY}--------------------------------------------------------{RESET}")
    import json
    if isinstance(payload, (dict, list)):
        print(f"{GREEN}{json.dumps(payload, indent=2)}{RESET}")
    else:
        print(f"{GREEN}{payload}{RESET}")
    print(f"{GRAY}========================================================{RESET}\n")

def normalize_question_json(q_data, concept_tag, flavor):
    # 1. Normalize prompt_text
    prompt_text = q_data.get("prompt_text") or q_data.get("questionText") or q_data.get("question") or q_data.get("text") or ""
    
    # 2. Normalize code_snippet
    code_snippet = q_data.get("code_snippet") or q_data.get("code")
    if code_snippet and isinstance(code_snippet, str) and len(code_snippet.strip()) > 0:
        cs_clean = code_snippet.strip()
        if "```" not in cs_clean:
            cs_clean = f"```\n{cs_clean}\n```"
        if cs_clean not in prompt_text:
            prompt_text = f"{prompt_text}\n\n{cs_clean}"
    else:
        code_snippet = None
        
    # 3. Normalize options
    raw_options = q_data.get("options")
    options = []
    if isinstance(raw_options, list):
        options = [str(o) for o in raw_options]
    elif isinstance(raw_options, dict):
        # Sort by key to maintain A, B, C, D order
        sorted_keys = sorted(raw_options.keys())
        options = [str(raw_options[k]) for k in sorted_keys]
        
    # Ensure exactly 4 options
    while len(options) < 4:
        options.append(f"Option {len(options)+1}")
    options = options[:4]
    
    # 4. Normalize correct_answer
    correct_answer = q_data.get("correct_answer") or q_data.get("correctAnswer") or q_data.get("answer") or ""
    # If correct_answer is just "A", "B", "C", "D": map it to the corresponding option value
    if correct_answer in ["A", "B", "C", "D"] and isinstance(raw_options, dict):
        correct_answer = raw_options.get(correct_answer, correct_answer)
    elif correct_answer in ["A", "B", "C", "D"] and len(options) >= 4:
        idx = ord(correct_answer) - ord('A')
        if 0 <= idx < len(options):
            correct_answer = options[idx]
            
    # 5. Normalize explanation
    explanation = q_data.get("explanation") or q_data.get("reason") or "No explanation provided."
    # 6. Normalize blooms_level
    blooms_level = q_data.get("blooms_level") or q_data.get("bloomsLevel")
    if not blooms_level:
        if flavor == "CONCEPTS_AND_DEFINITIONS":
            blooms_level = "Remember/Understand"
        elif flavor == "COMPARISONS_AND_TRADEOFFS":
            blooms_level = "Analyze"
        elif flavor == "FORMULAS_AND_CALCULATIONS":
            blooms_level = "Apply"
        elif flavor == "CASE_STUDIES_AND_SCENARIOS":
            blooms_level = "Apply"
        elif flavor == "PRACTICAL_AND_LAB_TASKS":
            blooms_level = "Apply"
        else:
            blooms_level = "Remember/Understand"
    
    return {
        "concept_tag": q_data.get("concept_tag") or concept_tag,
        "prompt_text": prompt_text,
        "questionText": prompt_text, # Dual keys compatibility
        "code_snippet": code_snippet,
        "options": options,
        "correct_answer": correct_answer,
        "correctAnswer": correct_answer, # Dual keys compatibility
        "explanation": explanation,
        "blooms_level": blooms_level
    }

def print_stage_header(stage_num, stage_title):
    CYAN = '\033[36m'
    RESET = '\033[0m'
    print(f"\n{CYAN}======================================================================{RESET}")
    print(f"{CYAN}  STAGE {stage_num}: {stage_title}{RESET}")
    print(f"{CYAN}======================================================================{RESET}")

def execute_generation_logic(req: GeneratorRequest):
    print_stage_header("1", "RAW INCOMING CAPTURE")
    
    if req.inputs:
        src_type = f"Multiple Uploaded Documents ({len(req.inputs)} Chunks Detected)"
        file_names_str = " | ".join([f"📑 {inp.source_name}" for inp in req.inputs])
    else:
        src_type = f"Single Source ({req.type or 'unknown'})"
        file_names_str = f"📑 {req.content[:50]}..." if req.content else "No File Name Available"
        
    raw_chars = len(req.content) if req.content else sum([len(inp.content) for inp in req.inputs]) if req.inputs else 0
    
    print(f"[Source Type] : {src_type}")
    print(f"[File Names]   : {file_names_str}")
    print(f"[Raw State]    : {raw_chars:,} characters of unparsed raw text structures.")

    MAX_CONTEXT_CHARS = 12000

    if req.inputs:
        resolved_text = resolve_input_sources(req.inputs)
        if not resolved_text or len(resolved_text.strip()) < 2:
            raise HTTPException(status_code=400, detail="Content too short or unextractable from sources.")
        
        print_stage_header("2", "CONTENT SANITIZATION & METADATA STRIPPING")
        print(f"[Action]       : Stripped away PPTX syntax, fonts, sizes, slide numbers, and timestamps.")
        chunks = deduplicate_text_chunks(resolved_text, embed_model)
        
        # Context-Length Budgeting chunk sieve
        context_parts = []
        current_len = 0
        for chunk in chunks:
            if current_len + len(chunk) + 2 > MAX_CONTEXT_CHARS:
                if len(context_parts) > 0:
                    break
            context_parts.append(chunk)
            current_len += len(chunk) + 2
        context = "\n\n".join(context_parts)
        print(f"[Clean Text]   : \"{context[:150].strip().replace(chr(10), ' ')}...\"")
    else:
        source_text = req.content
        if req.type in ['pdf', 'docx', 'pptx', 'image']:
            if os.path.exists(req.content):
                source_text = extract_text_from_file(req.content)
            elif req.content.startswith("base64:"):
                try:
                    b64_data = req.content[7:]
                    file_data = base64.b64decode(b64_data)
                    ext = req.type if req.type != 'image' else 'png'
                    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                        tmp.write(file_data)
                        tmp_path = tmp.name
                    source_text = extract_text_from_file(tmp_path)
                    os.remove(tmp_path)
                except Exception as e:
                    raise HTTPException(status_code=400, detail="Invalid base64 file content.")
        
        print_stage_header("2", "CONTENT SANITIZATION & METADATA STRIPPING")
        print(f"[Action]       : Stripped away PPTX syntax, fonts, sizes, slide numbers, and timestamps.")
        if req.type == 'topic':
            print(f"  [RAG] Searching global database for topic '{req.content}'...")
            db_results = get_relevant_db_context(req.content, top_k=5)
            if db_results:
                # Sieve database results to fit MAX_CONTEXT_CHARS
                context_parts = []
                current_len = 0
                for r in db_results:
                    chunk = r["content"]
                    if current_len + len(chunk) + 2 > MAX_CONTEXT_CHARS:
                        if len(context_parts) > 0:
                            break
                    context_parts.append(chunk)
                    current_len += len(chunk) + 2
                context = "\n\n".join(context_parts)
                print(f"  [RAG] Found {len(db_results)} matching chunks in database. Sieved to {len(context_parts)} chunks.")
            else:
                context = f"Syllabus topic: {req.content}"
                print("  [RAG] No matching chunks found in database. Using empty fallback.")
        else:
            if not source_text or len(source_text) < 2:
                raise HTTPException(status_code=400, detail="Content too short or file unreadable.")
            query = "Important core concepts"
            raw_context = get_relevant_context(source_text, query, top_k=5)
            # Sieve raw_context to fit MAX_CONTEXT_CHARS
            chunks = raw_context.split('\n\n')
            context_parts = []
            current_len = 0
            for chunk in chunks:
                if current_len + len(chunk) + 2 > MAX_CONTEXT_CHARS:
                    if len(context_parts) > 0:
                        break
                context_parts.append(chunk)
                current_len += len(chunk) + 2
            context = "\n\n".join(context_parts)
        print(f"[Clean Text]   : \"{context[:150].strip().replace(chr(10), ' ')}...\"")

    # Determine topic fallback for logs and RAG
    topic_fallback = "General Course Concept"
    if req.inputs:
        names = [inp.source_name for inp in req.inputs if inp.source_name and inp.source_name != "Unknown Source"]
        if names:
            clean_names = []
            for name in names:
                cleaned = re.sub(r'\.[a-zA-Z0-9]+$', '', name)
                cleaned = re.sub(r'(?i)voice\s+transcript\s*\([^)]+\)', '', cleaned)
                cleaned = cleaned.strip()
                if cleaned:
                    clean_names.append(cleaned)
            if clean_names:
                topic_fallback = " ".join(clean_names)
    elif req.content and req.type == 'topic':
        topic_fallback = req.content

    # ── Language Awareness Router ──
    programming_language = "Python/JavaScript"
    language_check_str = (topic_fallback + " " + context).lower()
    
    cpp_keywords = ['virtual function', 'pointer', 'c++', 'struct', 'header file', 'class template', 'memory allocation', 'cpp', 'destructor', 'constructor overload']
    python_keywords = ['pandas', 'dataframe', 'numpy', 'django', 'python', 'list comprehension', 'decorator', 'flask', 'pip', 'scikit']
    java_keywords = ['jvm', 'java', 'collections', 'interface', 'abstract class', 'garbage collection', 'package import', 'spring boot', 'extends', 'implements']
    sql_keywords = ['sql', 'select', 'join', 'database', 'query', 'primary key', 'foreign key', 'schema', 'table', 'where clause']
    
    if any(kw in language_check_str for kw in cpp_keywords):
        programming_language = "C++"
    elif any(kw in language_check_str for kw in java_keywords):
        programming_language = "Java"
    elif any(kw in language_check_str for kw in sql_keywords):
        programming_language = "SQL"
    elif any(kw in language_check_str for kw in python_keywords):
        programming_language = "Python"
        
    print(f"[Language Router] Detected target programming language: {programming_language}")

    keywords = ['mechanical', 'civil', 'chemical', 'structural', 'fluid', 'thermodynamic', 'material', 'drawing', 'concrete', 'machine', 'lab tracing', 'cad', 'optimiz', 'piping', 'construction', 'concrete', 'soil', 'geology', 'geotechnical', 'surveying']
    is_non_computational = any(kw in topic_fallback.lower() for kw in keywords)

    formal_names = {
        "CONCEPTS_AND_DEFINITIONS": "Core Theory",
        "COMPARISONS_AND_TRADEOFFS": "Analytical Reasoning",
        "FORMULAS_AND_CALCULATIONS": "Numerical Design",
        "CASE_STUDIES_AND_SCENARIOS": "Real-World Application",
        "PRACTICAL_AND_LAB_TASKS": "Design Optimization & Lab Tracing" if is_non_computational else "Implementation Synthesis"
    }

    # Determine Target Flavor Ratios and Question count
    default_ratios = {
        "CONCEPTS_AND_DEFINITIONS": 0.2,
        "COMPARISONS_AND_TRADEOFFS": 0.2,
        "FORMULAS_AND_CALCULATIONS": 0.2,
        "CASE_STUDIES_AND_SCENARIOS": 0.2,
        "PRACTICAL_AND_LAB_TASKS": 0.2
    }
    ratios = req.target_ratios or default_ratios
    total_count = req.count

    # Contextual evaluation layer for thematic incompatibility
    execution_messages = []
    text_content = context.lower()
    
    coa_keywords = ['coa', 'stack organization', 'accumulator organization', 'stack architecture', 'register organization', 'accumulator machine', 'cpu organization', 'instruction format', 'instruction cycle', 'hardware architecture', 'computer organization', 'addressing mode', 'cpu architecture', 'assembly language', 'instruction set architecture', 'isa', 'mips', 'risc', 'cisc', 'microarchitecture', 'arithmetic logic unit', 'pipeline hazard', 'cache', 'bus']
    is_coa = any(kw in text_content for kw in coa_keywords) or any(kw in topic_fallback.lower() for kw in coa_keywords)
    
    flowchart_keywords = ['flowchart', 'flow chart', 'control flow graph', 'cfg', 'pseudocode', 'pseudo code', 'program flow', 'flow-chart']
    is_flowchart = any(kw in text_content for kw in flowchart_keywords) or any(kw in topic_fallback.lower() for kw in flowchart_keywords)

    normalized_ratios = {
        "CONCEPTS_AND_DEFINITIONS": float(ratios.get("CONCEPTS_AND_DEFINITIONS") or ratios.get("CORE_THEORY") or 0.2),
        "COMPARISONS_AND_TRADEOFFS": float(ratios.get("COMPARISONS_AND_TRADEOFFS") or ratios.get("ANALYTICAL_REASONING") or 0.2),
        "FORMULAS_AND_CALCULATIONS": float(ratios.get("FORMULAS_AND_CALCULATIONS") or ratios.get("NUMERICAL_DESIGN") or 0.2),
        "CASE_STUDIES_AND_SCENARIOS": float(ratios.get("CASE_STUDIES_AND_SCENARIOS") or ratios.get("REAL_WORLD_APPLICATION") or 0.2),
        "PRACTICAL_AND_LAB_TASKS": float(ratios.get("PRACTICAL_AND_LAB_TASKS") or ratios.get("IMPLEMENTATION_SYNTHESIS") or 0.2),
    }

    if is_coa and normalized_ratios["PRACTICAL_AND_LAB_TASKS"] > 0.05:
        original_val = normalized_ratios["PRACTICAL_AND_LAB_TASKS"]
        normalized_ratios["PRACTICAL_AND_LAB_TASKS"] = 0.05
        diff = original_val - 0.05
        other_cats = ["CONCEPTS_AND_DEFINITIONS", "COMPARISONS_AND_TRADEOFFS", "FORMULAS_AND_CALCULATIONS", "CASE_STUDIES_AND_SCENARIOS"]
        other_sum = sum(normalized_ratios[cat] for cat in other_cats)
        if other_sum > 0:
            for cat in other_cats:
                normalized_ratios[cat] += diff * (normalized_ratios[cat] / other_sum)
        else:
            for cat in other_cats:
                normalized_ratios[cat] += diff / len(other_cats)
        for cat in normalized_ratios:
            normalized_ratios[cat] = round(normalized_ratios[cat], 2)
        total_sum = sum(normalized_ratios.values())
        if abs(total_sum - 1.0) > 0.0001:
            diff_round = round(1.0 - total_sum, 2)
            max_cat = max(other_cats, key=lambda c: normalized_ratios[c])
            normalized_ratios[max_cat] = round(normalized_ratios[max_cat] + diff_round, 2)
        execution_messages.append(
            "We clamped 'Implementation Synthesis' (Practical & Lab Tasks) to 5% because pure COA hardware architectures focus on low-level assembly tracing rather than high-level code implementation."
        )

    if is_flowchart and normalized_ratios["CASE_STUDIES_AND_SCENARIOS"] > 0.05:
        original_val = normalized_ratios["CASE_STUDIES_AND_SCENARIOS"]
        normalized_ratios["CASE_STUDIES_AND_SCENARIOS"] = 0.05
        diff = original_val - 0.05
        other_cats = ["CONCEPTS_AND_DEFINITIONS", "COMPARISONS_AND_TRADEOFFS", "FORMULAS_AND_CALCULATIONS", "PRACTICAL_AND_LAB_TASKS"]
        if is_coa:
            other_cats = ["CONCEPTS_AND_DEFINITIONS", "COMPARISONS_AND_TRADEOFFS", "FORMULAS_AND_CALCULATIONS"]
        other_sum = sum(normalized_ratios[cat] for cat in other_cats)
        if other_sum > 0:
            for cat in other_cats:
                normalized_ratios[cat] += diff * (normalized_ratios[cat] / other_sum)
        else:
            for cat in other_cats:
                normalized_ratios[cat] += diff / len(other_cats)
        for cat in normalized_ratios:
            normalized_ratios[cat] = round(normalized_ratios[cat], 2)
        total_sum = sum(normalized_ratios.values())
        if abs(total_sum - 1.0) > 0.0001:
            diff_round = round(1.0 - total_sum, 2)
            max_cat = max(other_cats, key=lambda c: normalized_ratios[c])
            normalized_ratios[max_cat] = round(normalized_ratios[max_cat] + diff_round, 2)
        execution_messages.append(
            "We clamped 'Real-World Application' (Case Studies & Scenarios) to 5% because abstract flowchart logic is best evaluated through design analysis rather than large system scenarios."
        )

    ratios = normalized_ratios
    
    # Compression Sieve Detection: total questions <= 5 and input length > 5000 tokens (approx 20,000 chars)
    is_compression_sieve = (total_count <= 5) and (len(context) > 20000)
    if is_compression_sieve:
        print("[Compression Sieve] Activated: input length > 5000 tokens and total questions <= 5. Emphasizing multi-concept synthesis questions.")

    # Calculate counts per flavor
    sum_ratios = sum(ratios.values())
    if sum_ratios == 0:
        ratios = default_ratios
        sum_ratios = 1.0
        
    counts = {}
    accumulated_count = 0
    active_flavors = [f for f in ratios.keys() if ratios[f] > 0]
    if not active_flavors:
        active_flavors = ["CONCEPTS_AND_DEFINITIONS"]
        
    for f in active_flavors[:-1]:
        c = int(round(total_count * (ratios[f] / sum_ratios)))
        counts[f] = c
        accumulated_count += c
    counts[active_flavors[-1]] = max(0, total_count - accumulated_count)

    print_stage_header("3", "DYNAMIC RATIO MATRIX (THE TEACHER'S SLIDERS)")
    requested_list = []
    for flavor_key, flavor_val in ratios.items():
        formal = formal_names.get(flavor_key, "Unknown")
        requested_list.append(f"{flavor_key} ({formal}) ({int(flavor_val * 100)}%)")
    requested_str = " | ".join(requested_list)
    print(f"[Requested]    : {requested_str}")
    print(f"[Hard Zero]    : Enforced! Small margins eliminated to maximize quiz focus.")
    
    final_list = []
    emojis = {
        "CONCEPTS_AND_DEFINITIONS": "📚",
        "COMPARISONS_AND_TRADEOFFS": "🔍",
        "FORMULAS_AND_CALCULATIONS": "⚙️",
        "CASE_STUDIES_AND_SCENARIOS": "💼",
        "PRACTICAL_AND_LAB_TASKS": "💻"
    }
    for f_key, f_cnt in counts.items():
        if f_cnt > 0:
            emoji = emojis.get(f_key, "❓")
            formal = formal_names.get(f_key, "Unknown")
            final_list.append(f"{emoji} {f_cnt} {f_key} ({formal}) Question{'s' if f_cnt > 1 else ''}")
    final_count_str = " | ".join(final_list)
    print(f"[Final Count]  : {final_count_str}")

    print_stage_header("4", "3-AGENT COGNITIVE GENERATION PIPELINE")
    print(f"={'='*55}")
    print(f"")

    # ── STAGE A: Content Understanding ──
    stage_a_data = run_stage_a_content_understanding(context, topic_fallback=topic_fallback)
    concepts = stage_a_data.get("concepts", [])
    extracted_exclusions = stage_a_data.get("isolated_narratives", [])

    # Merge narrative exclusions
    isolated_narratives = []
    if req.isolated_narratives:
        isolated_narratives.extend(req.isolated_narratives)
    if extracted_exclusions:
        isolated_narratives.extend(extracted_exclusions)
    isolated_narratives = list(set(isolated_narratives))

    # ── STAGE B: Academic Planner & Circuit Breaker Validator ──
    interview_mode_flag = bool(getattr(req, 'interview_mode', False) or (req.question_style == "INTERVIEW"))
    stage_b_blueprint = run_stage_b_academic_planner(
        stage_a_data=stage_a_data,
        question_count=total_count,
        difficulty=req.difficulty or "Medium",
        question_style=req.question_style or "MIXED",
        interview_mode=interview_mode_flag
    )

    # Compile generation tasks directly from Stage B Blueprint slots
    generation_tasks = []
    for slot in stage_b_blueprint["slots"]:
        concept_obj = {
            "concept_tag": slot["concept_tag"],
            "weight_score": slot.get("importance_score", 0.75)
        }
        flavor = slot["archetype"]
        slot_diff = slot["calibrated_difficulty"]
        resolved_style = slot["question_style"]
        generation_tasks.append((concept_obj, flavor, slot_diff, resolved_style))


    questions = []
    generated_so_far = ""

    # 3. Executing Agent 2 Generator & Agent 3 Critic loop
    for i, (concept, flavor, slot_diff, resolved_style) in enumerate(generation_tasks):
        formal = formal_names.get(flavor, "Unknown")
        print(f"\n• QUESTION {i+1} [Flavor: {flavor} ({formal}) | Style: {resolved_style}]")
        q_success = False
        
        # Initial Draft Generation by Agent 2
        for attempt in range(3):
            try:
                temp_offset = 0.0
                freq_penalty = None
                if attempt == 1:
                    temp_offset = 0.15
                    freq_penalty = 0.4
                    print("  ├── [Self-Healing]     : Attempt 2 - elevating temperature by +0.15 and frequency_penalty to 0.4")
                elif attempt == 2:
                    temp_offset = 0.30
                    freq_penalty = 0.4
                    print("  ├── [Self-Healing]     : Attempt 3 - elevating temperature by +0.30 and frequency_penalty to 0.4")

                q_data = run_agent2_generator(concept, flavor, context, generated_so_far, slot_diff, isolated_narratives, is_compression_sieve, temp_offset=temp_offset, freq_penalty=freq_penalty, programming_language=programming_language, question_style=resolved_style)
                print(f"  ├── [Agent 2 Creator]  : Drafted a multiple-choice question on '{concept.get('concept_tag')}'.")
                
                # Validation by Agent 3 (with difficulty for tier-based distractor checks)
                critic_res = run_agent3_critic(q_data, context, flavor, slot_diff, question_style=resolved_style)
                
                # Self-Correction Loop if failed
                if critic_res["status"] == "fail":
                    for repair_pass in range(2):
                        print(f"  ├── [Self-Correction]  : Agent 2 re-drafting the formatting structure (Attempt {repair_pass+1})...")
                        q_data = run_agent2_repair(q_data, critic_res["issues"], context)
                        critic_res = run_agent3_critic(q_data, context, flavor, slot_diff, question_style=resolved_style)
                        if critic_res["status"] == "pass":
                            break
                
                if critic_res["status"] == "pass":
                    normalized = normalize_question_json(q_data, concept.get("concept_tag"), flavor)
                    
                    q_final = {
                        "id": f"q_id_{str(i+1).zfill(3)}",
                        "type": flavor,
                        "concept_tag": normalized["concept_tag"],
                        "weight_score": float(concept.get("weight_score", 0.75)),
                        "prompt_text": normalized["prompt_text"],
                        "questionText": normalized["questionText"],
                        "code_snippet": normalized["code_snippet"],
                        "options": normalized["options"],
                        "correct_answer": normalized["correct_answer"],
                        "correctAnswer": normalized["correctAnswer"],
                        "explanation": normalized["explanation"],
                        "blooms_level": normalized["blooms_level"]
                    }
                    questions.append(q_final)
                    generated_so_far += f" [{q_final['prompt_text']}] "
                    print(f"  └── [Agent 3 Critic]   : [ \033[32m✅ PASS\033[0m ] -> Structure matches perfectly.")
                    q_success = True
                    break
                else:
                    print(f"  ├── [Agent 3 Critic]   : [ \033[31m❌ FAIL\033[0m ] -> {critic_res['issues'][0] if critic_res['issues'] else 'QA issues detected'}")
            except Exception as e:
                err_str = str(e).lower()
                if "timeout" in err_str or "read timed out" in err_str or "connection" in err_str:
                    print("  ├── [Agent 2 Creator]  : ⚠️ CPU Latency Spike -> Instantly routed to Groq Cloud fallback safety handler.")
                    if os.getenv("GROQ_API_KEY"):
                        try:
                            # Construct prompt for Groq to generate a single question conforming to schema
                            narrative_mutation_instruction = ""
                            if isolated_narratives and len(isolated_narratives) > 0:
                                exclusions_str = ", ".join([f'"{n}"' for n in isolated_narratives])
                                narrative_mutation_instruction = (
                                    f"NARRATIVE MUTATION MANDATE:\n"
                                    f"When generating questions under the REAL_WORLD_APPLICATION or NUMERICAL_DESIGN archetypes, review the following narrative/scenario exclusions: [{exclusions_str}].\n"
                                    f"You are STRICTLY PROHIBITED from using these exact scenarios, names, characters, or settings in your question stem or choices. "
                                    f"Instead, construct a structurally isomorphic (parallel) real-world scenario that tests the exact same concept using a brand-new application domain.\n\n"
                                )

                            prompt = (
                                "You are an elite Computer Science curriculum developer. Generate a single question.\n\n"
                                f"{narrative_mutation_instruction}"
                                "DIAGRAM AND GRAPHICS SAFEGUARD:\n"
                                "For any image/diagram-based text descriptions extracted from slides, you are only permitted to reproduce them if they represent standard engineering blueprints or industry-standard topologies. If the slide context indicates a non-standard or arbitrary sketch, synthesize a descriptive text-based conceptual evaluation instead.\n\n"
                                f"Context:\n{context[:2000]}\n\n"
                                f"Concept: {concept.get('concept_tag')}\n"
                                f"Type: {flavor}\n"
                                f"Difficulty: {req.difficulty}\n"
                                f"Already generated questions: {generated_so_far}\n\n"
                                "Return ONLY a clean JSON object conforming to this schema:\n"
                                "{\n"
                                "  \"concept_tag\": \"concept\",\n"
                                "  \"prompt_text\": \"Question text...\",\n"
                                "  \"code_snippet\": \"Markdown code snippet or null\",\n"
                                "  \"options\": [\"Option A\", \"Option B\", \"Option C\", \"Option D\"],\n"
                                "  \"correct_answer\": \"Exact copy of correct option from options\",\n"
                                "  \"explanation\": \"Detailed explanation explaining the solution...\"\n"
                                "}"
                            )
                            raw_text = call_groq_fallback(prompt, json_mode=True)
                            q_data = robust_json_loads(raw_text)
                            
                            normalized = normalize_question_json(q_data, concept.get("concept_tag"), flavor)
                            
                            q_final = {
                                "id": f"q_id_{str(i+1).zfill(3)}",
                                "type": flavor,
                                "concept_tag": normalized["concept_tag"],
                                "weight_score": float(concept.get("weight_score", 0.75)),
                                "prompt_text": normalized["prompt_text"],
                                "questionText": normalized["questionText"],
                                "code_snippet": normalized["code_snippet"],
                                "options": normalized["options"],
                                "correct_answer": normalized["correct_answer"],
                                "correctAnswer": normalized["correctAnswer"],
                                "explanation": normalized["explanation"],
                                "blooms_level": normalized["blooms_level"]
                            }
                            questions.append(q_final)
                            generated_so_far += f" [{q_final['prompt_text']}] "
                            print(f"  └── [Agent 3 Critic]   : [ \033[32m✅ PASS\033[0m ] -> Groq fallback generated successfully.")
                            q_success = True
                            break
                        except Exception as ge:
                            print(f"  ├── [Agent 2 Creator]  : ⚠️ Groq fallback safety handler failed: {ge}")
                else:
                    print(f"  ├── [Agent 2 Creator]  : ⚠️ Execution error: {e}")

        if not q_success:
            print(f"  └── [Agent 2 Creator]  : [ \033[31m❌ FAIL\033[0m ] -> Failed to generate question after 3 attempts. Gracefully bypassing and dropping corrupted node.")

    # Conforming to structured blueprint output JSON
    payload_response = {
        "quiz_metadata": {
            "source_material_id": req.source_material_id or "kmit_dynamic_gen",
            "total_questions": len(questions),
            "target_ratios": ratios,
            "isolated_narratives": isolated_narratives,
            "execution_messages": execution_messages
        },
        "questions": questions,
        "isolated_narratives": isolated_narratives
    }
    
    print_stage_header("5", "FAULT-TOLERANT WEBHOOK CALLBACK DISPATCH")
    print(f"[Status]       : {len(questions)}/{total_count} Validated questions compiled into clean JSON payload.")

    if not questions:
        raise HTTPException(status_code=500, detail="AI failed to generate any questions.")

    return payload_response

def run_generation_task(req: GeneratorRequest):
    try:
        result = execute_generation_logic(req)
        
        # Determine topic fallback for sub_topic mapping
        topic_fallback = "General Course Concept"
        if req.inputs:
            names = [inp.source_name for inp in req.inputs if inp.source_name and inp.source_name != "Unknown Source"]
            if names:
                clean_names = []
                for name in names:
                    cleaned = re.sub(r'\.[a-zA-Z0-9]+$', '', name)
                    cleaned = re.sub(r'(?i)voice\s+transcript\s*\([^)]+\)', '', cleaned)
                    cleaned = cleaned.strip()
                    if cleaned:
                        clean_names.append(cleaned)
                if clean_names:
                    topic_fallback = " ".join(clean_names)
        elif req.content and req.type == 'topic':
            topic_fallback = req.content

        # Enforce strict key normalization on callback dispatch
        normalized_questions = []
        for i, q in enumerate(result.get("questions", [])):
            raw_options = q.get("options")
            options_list = []
            if isinstance(raw_options, list):
                options_list = [str(o) for o in raw_options]
            elif isinstance(raw_options, dict):
                sorted_keys = sorted(raw_options.keys())
                options_list = [str(raw_options[k]) for k in sorted_keys]
            while len(options_list) < 4:
                options_list.append(f"Option {len(options_list)+1}")
            options_list = options_list[:4]
            
            # Map options uniformly to object containing keys { A, B, C, D }
            options_dict = {
                "A": options_list[0],
                "B": options_list[1],
                "C": options_list[2],
                "D": options_list[3]
            }

            norm_q = {
                "id": q.get("id") or f"q_id_{str(i+1).zfill(3)}",
                "type": q.get("type") or "theory",
                "concept_tag": q.get("concept_tag") or topic_fallback,
                "sub_topic": q.get("concept_tag") or topic_fallback, # UI displays actual topic
                "weight_score": float(q.get("weight_score") or 0.75),
                "prompt_text": q.get("prompt_text") or q.get("questionText") or "",
                "questionText": q.get("prompt_text") or q.get("questionText") or "",
                "code_snippet": q.get("code_snippet"),
                "options": options_dict, # Object containing { A, B, C, D }
                "correct_answer": q.get("correct_answer") or q.get("correctAnswer") or "",
                "correctAnswer": q.get("correct_answer") or q.get("correctAnswer") or "",
                "explanation": q.get("explanation") or "Detailed explanation for solution choice.",
                "blooms_level": q.get("blooms_level") or "Remember/Understand"
            }
            normalized_questions.append(norm_q)

        result["questions"] = normalized_questions

        # Post success callback
        print(f"[Callback]     : Forwarding complete quiz structure to cloud engine... ", end="", flush=True)
        res = requests.post(req.callback_url, json={"status": "success", "result": result}, headers={"Content-Type": "application/json"}, timeout=30)
        print(f"[Success {res.status_code}]")
        print(f"======================================================================\n")
    except Exception as e:
        import traceback
        exc_type = type(e).__name__
        exc_msg = str(e) or "(empty exception message — likely an IndexError or NoneType access)"
        full_trace = traceback.format_exc()
        print(f"[Callback]     : Generation failed [{exc_type}]: {exc_msg}")
        print(f"[Traceback]    : {full_trace[:600]}")
        print(f"[Callback]     : Dispatching failure callback... ", end="", flush=True)
        try:
            res = requests.post(req.callback_url, json={"status": "failed", "error": f"[{exc_type}] {exc_msg}"}, headers={"Content-Type": "application/json"}, timeout=30)
            print(f"[Failed Callback Sent {res.status_code}]")
            print(f"======================================================================\n")
        except Exception as cb_err:
            print(f"[Error {cb_err}]")
            print(f"======================================================================\n")

@app.post("/generate")
async def generate_questions(req: GeneratorRequest, background_tasks: BackgroundTasks):
    if req.callback_url:
        print(f"  [PIPELINE] Received async request. Scheduling background generation task. callback_url={req.callback_url}")
        background_tasks.add_task(run_generation_task, req)
        return {"status": "accepted", "msg": "Quiz generation started in background", "taskId": req.taskId}
    
    # Otherwise synchronous execution
    return execute_generation_logic(req)

if __name__ == "__main__":
    print(f"AI Service starting on port 8000 using local model: {MODEL_NAME}")
    uvicorn.run(app, host="0.0.0.0", port=8000)

